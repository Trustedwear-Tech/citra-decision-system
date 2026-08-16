# Copyright (c) 2026 Trustedwear Tech Private Limited (https://citra-ai.com)
# Author: Rohit Kumar Chandan
# SPDX-License-Identifier: BUSL-1.1
#
# Licensed under the Business Source License 1.1. Non-production use is granted;
# production use requires a commercial licence until the Change Date, after
# which this file converts to Apache-2.0. See LICENSE at the repository root.

"""Document extraction (PDF / Office / HTML) — runs locally inside the
sandbox using small pure-Python parsers. For scanned/image pages we
fall back to ``citra_toolkit.ocr.extract`` (Qwen3-VL-32B via the proxy).

No tesseract, no poppler, no torch. All extractors are optional imports
so an extractor missing in the image only disables that path, not the
whole module.
"""
from __future__ import annotations

import io
import logging
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from . import ocr as _ocr

logger = logging.getLogger(__name__)


@dataclass
class ExtractedDoc:
    text: str
    pages: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    source_kind: str = ""  # pdf | docx | pptx | xlsx | html | text


# ----------------------------------------------------------------------- PDF
def extract_pdf(data: bytes | str | os.PathLike, *, ocr: bool = False) -> ExtractedDoc:
    """Extract text + tables from a PDF.

    ``data`` may be raw bytes or a path. If a page yields no text and
    ``ocr=True``, we render+OCR via the vision proxy.
    """
    if isinstance(data, (bytes, bytearray)):
        raw = bytes(data)
    else:
        raw = Path(data).read_bytes()

    pages: list[str] = []
    tables: list[list[list[str]]] = []
    meta: dict[str, Any] = {}

    # Text via pdfminer.six
    try:
        from pdfminer.high_level import extract_text  # type: ignore
        full = extract_text(io.BytesIO(raw)) or ""
        # pdfminer joins pages with form-feed
        pages = [p for p in full.split("\f") if p is not None]
    except Exception as e:  # noqa: BLE001
        logger.warning("pdfminer.six failed: %s", e)
        pages = []

    # Metadata via pypdf
    try:
        from pypdf import PdfReader  # type: ignore
        reader = PdfReader(io.BytesIO(raw))
        meta = {
            "num_pages": len(reader.pages),
            "info": {k: str(v) for k, v in (reader.metadata or {}).items()} if reader.metadata else {},
        }
        if not pages:
            pages = []
            for p in reader.pages:
                try:
                    pages.append(p.extract_text() or "")
                except Exception:  # noqa: BLE001
                    pages.append("")
    except Exception as e:  # noqa: BLE001
        logger.debug("pypdf metadata failed: %s", e)

    # Tables via pdfplumber
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page in pdf.pages:
                try:
                    page_tables = page.extract_tables() or []
                except Exception:  # noqa: BLE001
                    page_tables = []
                for tbl in page_tables:
                    if tbl:
                        tables.append([[(c or "") for c in row] for row in tbl])
    except Exception as e:  # noqa: BLE001
        logger.debug("pdfplumber failed: %s", e)

    # OCR fallback for empty pages
    if ocr and any(not (p and p.strip()) for p in pages):
        # We cannot rasterize without poppler; ask vision proxy to read the whole
        # PDF (the proxy supports application/pdf payloads).
        try:
            res = _ocr.extract(raw, mode="layout", content_type="application/pdf",
                               filename="upload.pdf")
            if res.text:
                pages = [res.text] if not pages else [(p or res.text) for p in pages]
        except Exception as e:  # noqa: BLE001
            logger.warning("ocr fallback failed: %s", e)

    text = "\n\n".join(p.strip() for p in pages if p)
    return ExtractedDoc(text=text, pages=pages, tables=tables, metadata=meta, source_kind="pdf")


# --------------------------------------------------------------------- DOCX
def extract_docx(data: bytes | str | os.PathLike) -> ExtractedDoc:
    from docx import Document  # type: ignore
    if isinstance(data, (bytes, bytearray)):
        doc = Document(io.BytesIO(bytes(data)))
    else:
        doc = Document(str(data))
    paragraphs = [p.text for p in doc.paragraphs]
    tables: list[list[list[str]]] = []
    for t in doc.tables:
        tables.append([[cell.text for cell in row.cells] for row in t.rows])
    return ExtractedDoc(text="\n".join(paragraphs), pages=paragraphs,
                        tables=tables, source_kind="docx")


# --------------------------------------------------------------------- PPTX
def extract_pptx(data: bytes | str | os.PathLike) -> ExtractedDoc:
    from pptx import Presentation  # type: ignore
    if isinstance(data, (bytes, bytearray)):
        prs = Presentation(io.BytesIO(bytes(data)))
    else:
        prs = Presentation(str(data))
    pages: list[str] = []
    for slide in prs.slides:
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
        pages.append("\n".join(chunks))
    return ExtractedDoc(text="\n\n".join(pages), pages=pages, source_kind="pptx")


# --------------------------------------------------------------------- XLSX
def extract_xlsx(data: bytes | str | os.PathLike) -> ExtractedDoc:
    from openpyxl import load_workbook  # type: ignore
    if isinstance(data, (bytes, bytearray)):
        wb = load_workbook(io.BytesIO(bytes(data)), data_only=True, read_only=True)
    else:
        wb = load_workbook(str(data), data_only=True, read_only=True)
    tables: list[list[list[str]]] = []
    text_parts: list[str] = []
    for ws in wb.worksheets:
        sheet: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            sheet.append(["" if c is None else str(c) for c in row])
        tables.append(sheet)
        text_parts.append(f"# {ws.title}\n" + "\n".join("\t".join(r) for r in sheet))
    return ExtractedDoc(text="\n\n".join(text_parts), tables=tables, source_kind="xlsx")


# --------------------------------------------------------------------- HTML
def extract_html(data: bytes | str) -> ExtractedDoc:
    raw = data.decode("utf-8", errors="replace") if isinstance(data, (bytes, bytearray)) else str(data)
    text = ""
    try:
        import trafilatura  # type: ignore
        text = trafilatura.extract(raw, include_tables=True, include_links=False) or ""
    except Exception:  # noqa: BLE001
        pass
    if not text:
        try:
            from readability import Document as RDoc  # type: ignore
            from bs4 import BeautifulSoup  # type: ignore
            summary_html = RDoc(raw).summary(html_partial=True)
            text = BeautifulSoup(summary_html, "html.parser").get_text("\n", strip=True)
        except Exception:  # noqa: BLE001
            try:
                from bs4 import BeautifulSoup  # type: ignore
                text = BeautifulSoup(raw, "html.parser").get_text("\n", strip=True)
            except Exception:  # noqa: BLE001
                text = ""
    return ExtractedDoc(text=text, source_kind="html")


# ------------------------------------------------------------------- router
_PDF_MAGIC = b"%PDF-"
_OFFICE_ZIP_MAGIC = b"PK\x03\x04"


def extract(data: bytes | str | os.PathLike, *,
            hint_mime: str | None = None,
            hint_filename: str | None = None,
            ocr: bool = False) -> ExtractedDoc:
    """Auto-route to the right extractor based on magic bytes / hints."""
    if isinstance(data, (str, os.PathLike)):
        path = Path(data)
        suffix = path.suffix.lower()
        raw = path.read_bytes()
        hint_filename = hint_filename or path.name
    else:
        raw = bytes(data)
        suffix = (Path(hint_filename).suffix.lower() if hint_filename else "")

    mime = (hint_mime or "").lower()

    if mime == "application/pdf" or suffix == ".pdf" or raw.startswith(_PDF_MAGIC):
        return extract_pdf(raw, ocr=ocr)
    if suffix == ".docx" or mime.endswith("wordprocessingml.document"):
        return extract_docx(raw)
    if suffix == ".pptx" or mime.endswith("presentationml.presentation"):
        return extract_pptx(raw)
    if suffix == ".xlsx" or mime.endswith("spreadsheetml.sheet"):
        return extract_xlsx(raw)
    if (mime.startswith("text/html") or suffix in (".html", ".htm")
            or raw[:64].lstrip().lower().startswith(b"<!doctype html")
            or raw[:64].lstrip().lower().startswith(b"<html")):
        return extract_html(raw)
    if mime.startswith("text/") or suffix in (".txt", ".md", ".csv", ".log", ".json"):
        try:
            return ExtractedDoc(text=raw.decode("utf-8", errors="replace"), source_kind="text")
        except Exception:  # noqa: BLE001
            return ExtractedDoc(text="", source_kind="text")
    if mime.startswith("image/") and ocr:
        res = _ocr.extract(raw, mode="layout", content_type=mime or "image/png")
        return ExtractedDoc(text=res.text, source_kind="image")

    # Last-resort: try each extractor in order, swallow failures.
    for fn, kind in [(extract_pdf, "pdf"), (extract_html, "html")]:
        try:
            return fn(raw)  # type: ignore[arg-type]
        except Exception:  # noqa: BLE001
            continue
    return ExtractedDoc(text="", source_kind="unknown")


# ------------------------------------------------------------------- chunker
def chunk(text: str, *, target_tokens: int = 512, overlap: int = 64) -> list[str]:
    """Sentence-aware chunker. Approximates tokens as ``len(words) * 1.3``.

    Falls back to a fixed-size character chunker when ``pysbd`` is not
    available.
    """
    if not text or not text.strip():
        return []
    sentences: list[str]
    try:
        import pysbd  # type: ignore
        seg = pysbd.Segmenter(language="en", clean=False)
        sentences = [s for s in seg.segment(text) if s and s.strip()]
    except Exception:  # noqa: BLE001
        sentences = [s.strip() for s in text.split(".") if s.strip()]
    if not sentences:
        return []

    target_words = max(50, int(target_tokens / 1.3))
    overlap_words = max(0, int(overlap / 1.3))

    chunks: list[str] = []
    cur: list[str] = []
    cur_words = 0
    for s in sentences:
        w = len(s.split())
        if cur and cur_words + w > target_words:
            chunks.append(" ".join(cur).strip())
            # carry overlap
            if overlap_words and cur:
                tail: list[str] = []
                tail_w = 0
                for s2 in reversed(cur):
                    tail_w += len(s2.split())
                    tail.insert(0, s2)
                    if tail_w >= overlap_words:
                        break
                cur = list(tail)
                cur_words = sum(len(x.split()) for x in cur)
            else:
                cur = []
                cur_words = 0
        cur.append(s)
        cur_words += w
    if cur:
        chunks.append(" ".join(cur).strip())
    return [c for c in chunks if c]
