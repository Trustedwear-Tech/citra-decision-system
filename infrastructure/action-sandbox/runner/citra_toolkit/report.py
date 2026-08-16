# Copyright (c) 2026 Trustedwear Tech Private Limited (https://citra-ai.com)
# Author: Rohit Kumar Chandan
# SPDX-License-Identifier: BUSL-1.1
#
# Licensed under the Business Source License 1.1. Non-production use is granted;
# production use requires a commercial licence until the Change Date, after
# which this file converts to Apache-2.0. See LICENSE at the repository root.

"""Generate report artifacts inside the sandbox.

Outputs land in ``$CITRA_WORKSPACE_DIR/reports/<session_id>/`` (tmpfs).
The directory dies with the container, so each ``make_*`` helper
**auto-publishes** the artifact by default — it queues the file for
upload to Citra-Service object storage (S3 / MinIO) which returns a
presigned download URL the Action Chat UI surfaces to the user.

Pass ``auto_publish=False`` if you need to post-process the file before
publishing (e.g. embedding it into another report). Manual
``citra_toolkit.publish.publish()`` is idempotent, so calling it on a
path that was already auto-published is a safe no-op.

Heavy deps (weasyprint, openpyxl, matplotlib, pandas) are baked into the
sandbox image — we never ``pip install`` at runtime. A Chart.js renderer
(Node) is baked in too; ``make_chartjs_png`` shells out to it for
web-grade charts, while ``make_chart_png`` (matplotlib) covers in-process
analytical plots.
"""
from __future__ import annotations

import base64
import csv
import io
import os
import re
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any, Iterable, Mapping

from . import publish as _publish


_SLUG_RE = re.compile(r"[^a-zA-Z0-9._-]+")


def _maybe_publish(
    path: Path,
    *,
    kind: str,
    title: str,
    summary: str,
    auto_publish: bool,
) -> None:
    if not auto_publish:
        return
    try:
        _publish.publish(path, kind=kind, title=title or path.stem, summary=summary)  # type: ignore[arg-type]
    except Exception:  # noqa: BLE001
        # Never let an outbox-write failure break artifact generation.
        # The agent can still call publish() manually if it cares.
        pass


def _reports_dir() -> Path:
    sid = os.getenv("CITRA_SESSION_ID", "unknown")
    workspace = Path(os.getenv("CITRA_WORKSPACE_DIR", "/workspace"))
    out = workspace / "reports" / sid
    out.mkdir(parents=True, exist_ok=True)
    return out


def _stamped_path(slug: str, ext: str) -> Path:
    safe = _SLUG_RE.sub("-", slug.strip().lower())[:48] or "report"
    ts = datetime.now(tz=timezone.utc).strftime("%Y%m%d-%H%M%S")
    return _reports_dir() / f"{ts}-{safe}.{ext.lstrip('.')}"


# ------------------------------- SVG → PNG --------------------------------
#
# Both make_pptx (python-pptx) and make_chart_png consumers expect raster
# bytes; python-pptx in particular cannot embed raw SVG. WeasyPrint (the
# PDF path) renders inline ``<svg>`` natively, so SVG bytes never need to
# be rasterised for ``make_pdf`` — pass the SVG as inline HTML and it
# just works.
#
# This helper is exposed at module level so skills can call it directly
# when they want to control the raster size, but ``coerce_image`` inside
# make_pptx also auto-detects SVG content and rasterises on the fly. So
# slides can pass an SVG string in the ``image`` field with no ceremony.


def _looks_like_svg(value: Any) -> bool:
    """Return True if value is an SVG markup string or bytes."""
    if isinstance(value, (bytes, bytearray)):
        head = bytes(value)[:512].lstrip()
        return head.startswith(b"<svg") or (
            head.startswith(b"<?xml") and b"<svg" in head[:512]
        )
    if isinstance(value, str):
        head = value.lstrip()[:512]
        return head.startswith("<svg") or (
            head.startswith("<?xml") and "<svg" in head
        )
    return False


def svg_to_png(
    svg: str | bytes,
    *,
    width: int | None = None,
    height: int | None = None,
    scale: float = 2.0,
) -> bytes:
    """Rasterise an SVG string/bytes to PNG bytes.

    Use this when a slide spec or report path needs SVG as raster. For
    PPTX, the agent can also just put the SVG string in the ``image``
    field — coerce_image() will call this helper transparently. For
    PDF/HTML, prefer inline ``<svg>`` (WeasyPrint renders it natively).

    ``scale`` controls oversampling for crisp output (2x default).
    Explicit ``width``/``height`` override the SVG's own size.

    Backed by cairosvg, which uses libcairo (already installed in the
    sandbox image for WeasyPrint).
    """
    if isinstance(svg, str):
        svg_bytes = svg.encode("utf-8")
    else:
        svg_bytes = bytes(svg)

    try:
        from cairosvg import svg2png  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover — surfaced at runtime
        raise RuntimeError(
            "SVG rasterisation requires cairosvg. "
            "Install cairosvg==2.7.1 (libcairo is already in the sandbox)."
        ) from exc

    kwargs: dict[str, Any] = {"bytestring": svg_bytes}
    if width and height:
        kwargs["output_width"] = int(width)
        kwargs["output_height"] = int(height)
    elif scale and scale != 1.0:
        kwargs["scale"] = float(scale)

    result = svg2png(**kwargs)
    if not isinstance(result, (bytes, bytearray)):
        raise ValueError("cairosvg returned no PNG bytes")
    return bytes(result)


# ------------------------------- PDF ---------------------------------------
def make_pdf(
    html: str,
    *,
    title: str = "report",
    summary: str = "",
    auto_publish: bool = True,
) -> Path:
    """Render an HTML string to a single-file PDF using WeasyPrint.

    Auto-publishes the file by default (uploaded to S3/MinIO; presigned
    URL surfaces in the chat). Pass ``auto_publish=False`` to defer.
    """
    from weasyprint import HTML  # heavy import, defer

    out = _stamped_path(title, "pdf")
    HTML(string=html, base_url=str(_reports_dir())).write_pdf(str(out))
    _maybe_publish(out, kind="pdf", title=title, summary=summary, auto_publish=auto_publish)
    return out


# ---------------- Executive PDF templates ----------------------------------
# Mirrors the Citra-Service web pipeline's exec_pg_* family
# (`Citra-Service/printable/printable_templates.py`). Same 7 section types,
# same palette, same rhythm rules — agent passes a list of structured
# section dicts and ships a polished A4 report.

# Allowed exec PDF section types.
_PDF_EXEC_SECTIONS = (
    "cover",             # page 1 — dark navy, two-tone headline, pill tags
    "argument",          # workhorse — kicker + action title + bullets card
    "stat_grid",         # 2×2 KPI grid with accent bars
    "features_2x2",      # 2×2 capability grid with icon circles
    "industries_2x2",    # 2×2 industry grid with side-rules + check uses
    "sovereignty_dark",  # dark — 4 trust cards + light governance panel
    "closing_dark",      # last — dark, 2×2 numbered reason cards + CTA
)

# Shared CSS — palette + per-section component styles. Kept inline so the
# rendered HTML is self-contained (WeasyPrint can read no extra files).
_EXEC_PDF_CSS = """
@page { size: A4 portrait; margin: 0; }
* { box-sizing: border-box; }
body { margin: 0; font-family: "Calibri","Segoe UI","Helvetica Neue",sans-serif;
       color: #0F172A; font-size: 11pt; line-height: 1.5; }

/* Per-page frame: 210×297mm, padded, page-broken between sections.
   The 100vh / page-break trick gives each section its own A4 page. */
.pg { width: 210mm; min-height: 297mm; padding: 14mm 14mm 16mm 14mm;
      page-break-after: always; position: relative; overflow: hidden; }
.pg:last-child { page-break-after: avoid; }
.pg-light { background: #F8FAFC; color: #0F172A; }
.pg-dark  { background: #0B1020; color: #FFFFFF; }

/* Common typography */
.kicker { font-size: 10pt; font-weight: 700; letter-spacing: 3px;
          text-transform: uppercase; color: #2563EB; margin: 0 0 6mm 0; }
.kicker-cyan { color: #22D3EE; }
.kicker-purple { color: #8B5CF6; }
.kicker-cyan2 { color: #06B6D4; }
.action-title { font-size: 26pt; font-weight: 700; color: #0F172A;
                line-height: 1.15; margin: 0 0 4mm 0; }
.action-title.on-dark { color: #FFFFFF; }
.subhead { font-size: 12pt; color: #475569; line-height: 1.55;
           margin: 0 0 7mm 0; }
.subhead.on-dark { color: #CBD5E1; }

/* Cover */
.cover-brand { display: inline-block; padding: 4mm 6mm;
               border: 1px solid #475569; border-radius: 4px;
               font-size: 10pt; font-weight: 700; letter-spacing: 3px;
               text-transform: uppercase; color: #FFFFFF; margin-bottom: 50mm; }
.cover-kicker { font-size: 11pt; font-weight: 700; letter-spacing: 4px;
                text-transform: uppercase; color: #22D3EE; margin: 0 0 4mm 0; }
.cover-title-a { font-size: 44pt; font-weight: 700; line-height: 1.08;
                 color: #FFFFFF; margin: 0; }
.cover-title-b { font-size: 44pt; font-weight: 700; line-height: 1.08;
                 color: #22D3EE; margin: 0; }
.cover-subhead { font-size: 14pt; color: #CBD5E1; line-height: 1.55;
                 margin: 14mm 0 0 0; max-width: 165mm; }
.cover-pills { position: absolute; left: 14mm; right: 14mm; bottom: 18mm;
               display: flex; gap: 6mm; }
.cover-pill { flex: 1; background: #1E293B; color: #FFFFFF;
              border-radius: 8mm; padding: 3mm 5mm; text-align: center;
              font-size: 11pt; font-weight: 700; }

/* Argument (workhorse) */
.arg-card { background: #FFFFFF; border-radius: 4mm; padding: 9mm 10mm;
            box-shadow: 0 2mm 4mm rgba(15,23,42,0.06); margin-top: 2mm; }
.arg-section-heading { font-size: 16pt; font-weight: 700; color: #0F172A;
                       margin: 0 0 3mm 0; }
.arg-intro { font-size: 12pt; color: #334155; line-height: 1.65;
             margin: 0 0 6mm 0; }
.arg-bullets { list-style: none; margin: 0; padding: 0; }
.arg-bullets li { position: relative; padding-left: 8mm;
                  font-size: 12pt; color: #1F2937; line-height: 1.7;
                  margin: 0 0 3mm 0; }
.arg-bullets li::before { content: "•"; position: absolute; left: 0;
                          color: #2563EB; font-weight: 700; font-size: 14pt; }
.arg-takeaway { margin-top: 7mm; padding: 4mm 5mm;
                text-align: center; font-size: 11pt; font-weight: 700;
                letter-spacing: 2px; text-transform: uppercase;
                color: #2563EB; }

/* Stat grid */
.stat-grid { display: grid; grid-template-columns: 1fr 1fr;
             gap: 6mm; margin-top: 2mm; }
.stat-card { background: #FFFFFF; border-radius: 3mm; padding: 8mm 9mm;
             box-shadow: 0 1mm 3mm rgba(15,23,42,0.05);
             border-top: 1.5mm solid #2563EB; min-height: 50mm; }
.stat-card.accent-blue   { border-top-color: #2563EB; }
.stat-card.accent-cyan   { border-top-color: #06B6D4; }
.stat-card.accent-purple { border-top-color: #8B5CF6; }
.stat-card.accent-green  { border-top-color: #10B981; }
.stat-value { font-size: 36pt; font-weight: 700; line-height: 1.0;
              color: #0F172A; margin: 0; }
.stat-label { font-size: 11pt; color: #475569; line-height: 1.45;
              margin: 4mm 0 0 0; }
.stat-detail { margin-top: 8mm; }
.stat-detail-heading { font-size: 14pt; font-weight: 700; color: #0F172A;
                       margin: 0 0 4mm 0; }
.stat-detail-body { font-size: 11pt; color: #334155; line-height: 1.7;
                    margin: 0; }

/* Features 2×2 */
.feat-grid { display: grid; grid-template-columns: 1fr 1fr;
             gap: 6mm; margin-top: 2mm; }
.feat-card { background: #FFFFFF; border-radius: 3mm; padding: 8mm;
             box-shadow: 0 1mm 3mm rgba(15,23,42,0.05); min-height: 60mm; }
.feat-head { display: flex; align-items: center; gap: 5mm; margin-bottom: 5mm; }
.feat-icon { width: 14mm; height: 14mm; border-radius: 50%;
             color: #FFFFFF; font-size: 14pt; font-weight: 700;
             display: inline-flex; align-items: center; justify-content: center; }
.feat-icon.icon-blue   { background: #2563EB; }
.feat-icon.icon-cyan   { background: #06B6D4; }
.feat-icon.icon-purple { background: #8B5CF6; }
.feat-icon.icon-green  { background: #10B981; }
.feat-title { font-size: 14pt; font-weight: 700; color: #0F172A; margin: 0; }
.feat-body { font-size: 11pt; color: #475569; line-height: 1.6; margin: 0; }
.feat-banner { background: #0B1020; color: #22D3EE; border-radius: 2mm;
               padding: 4mm; text-align: center; font-size: 11pt;
               font-weight: 700; letter-spacing: 2px;
               text-transform: uppercase; margin-top: 6mm; }

/* Industries 2×2 */
.ind-grid { display: grid; grid-template-columns: 1fr 1fr;
            gap: 6mm; margin-top: 2mm; }
.ind-card { background: #FFFFFF; border-radius: 3mm; padding: 7mm 7mm 7mm 9mm;
            box-shadow: 0 1mm 3mm rgba(15,23,42,0.05); min-height: 60mm;
            border-left: 1.5mm solid #2563EB; }
.ind-card.rule-blue   { border-left-color: #2563EB; }
.ind-card.rule-green  { border-left-color: #10B981; }
.ind-card.rule-orange { border-left-color: #F59E0B; }
.ind-card.rule-purple { border-left-color: #8B5CF6; }
.ind-head { display: flex; align-items: center; gap: 5mm; margin-bottom: 5mm; }
.ind-icon { width: 10mm; height: 10mm; border-radius: 50%;
            color: #FFFFFF; font-size: 11pt; font-weight: 700;
            display: inline-flex; align-items: center; justify-content: center; }
.ind-icon.icon-blue   { background: #2563EB; }
.ind-icon.icon-green  { background: #10B981; }
.ind-icon.icon-orange { background: #F59E0B; }
.ind-icon.icon-purple { background: #8B5CF6; }
.ind-name { font-size: 15pt; font-weight: 700; color: #0F172A; margin: 0; }
.ind-uses { columns: 2; column-gap: 5mm; margin: 0; padding: 0;
            list-style: none; }
.ind-uses li { position: relative; padding-left: 6mm;
               font-size: 10.5pt; color: #334155; line-height: 1.75;
               break-inside: avoid; margin-bottom: 2mm; }
.ind-uses li::before { content: "✓"; position: absolute; left: 0;
                       color: #10B981; font-weight: 700; }

/* Sovereignty dark */
.sov-grid { display: grid; grid-template-columns: 1fr 1fr;
            gap: 5mm; margin-top: 2mm; }
.sov-card { background: #111827; border: 1px solid #1E293B;
            border-radius: 3mm; padding: 7mm 8mm; min-height: 52mm; }
.sov-card-title { font-size: 18pt; font-weight: 700; color: #22D3EE;
                  margin: 0 0 4mm 0; }
.sov-card-body { font-size: 11pt; color: #FFFFFF; line-height: 1.6; margin: 0; }
.sov-gov { background: #F8FAFC; color: #0F172A; border-radius: 3mm;
           padding: 7mm 9mm; margin-top: 6mm; }
.sov-gov-heading { font-size: 11pt; font-weight: 700; letter-spacing: 3px;
                   text-transform: uppercase; color: #2563EB;
                   margin: 0 0 5mm 0; }
.sov-gov-grid { display: grid; grid-template-columns: 1fr 1fr;
                gap: 5mm 8mm; }
.sov-gov-item { display: grid; grid-template-columns: 8mm 1fr; gap: 3mm;
                align-items: start; }
.sov-gov-check { color: #10B981; font-weight: 700; font-size: 12pt; }
.sov-gov-title { font-size: 11pt; font-weight: 700; color: #0F172A;
                 margin: 0 0 1mm 0; }
.sov-gov-body { font-size: 10pt; color: #475569; line-height: 1.5; margin: 0; }

/* Closing dark */
.close-grid { display: grid; grid-template-columns: 1fr 1fr;
              gap: 5mm; margin-top: 2mm; }
.close-card { background: #111827; border: 1px solid #1E293B;
              border-radius: 3mm; padding: 7mm 8mm; min-height: 50mm;
              display: grid; grid-template-columns: auto 1fr;
              gap: 6mm; align-items: start; }
.close-num { font-size: 32pt; font-weight: 700; color: #22D3EE;
             line-height: 1.0; }
.close-title { font-size: 14pt; font-weight: 700; color: #FFFFFF;
               margin: 0 0 2mm 0; }
.close-body { font-size: 11pt; color: #CBD5E1; line-height: 1.55; margin: 0; }
.close-cta { background: #22D3EE; color: #0B1020; border-radius: 3mm;
             padding: 5mm 6mm; margin-top: 8mm; text-align: center;
             font-size: 13pt; font-weight: 700; letter-spacing: 1px;
             text-transform: uppercase; }
"""


def _esc(value: Any) -> str:
    """HTML-escape any value (None → empty string)."""
    if value is None:
        return ""
    import html as _html
    return _html.escape(str(value))


def _pdf_render_cover(sec: Mapping[str, Any]) -> str:
    brand = _esc(sec.get("brand_chip"))
    kicker = _esc(sec.get("kicker"))
    title_a = _esc(sec.get("title_a") or sec.get("title") or "")
    title_b = _esc(sec.get("title_b") or sec.get("subtitle") or "")
    subhead = _esc(sec.get("subhead"))
    pills = [_esc(p) for p in (sec.get("pills") or [])][:3]
    brand_html = (
        f'<div class="cover-brand">{brand}</div>' if brand else ""
    )
    kicker_html = (
        f'<div class="cover-kicker">{kicker}</div>' if kicker else ""
    )
    pills_html = ""
    if pills:
        pill_items = "".join(
            f'<div class="cover-pill">{p}</div>' for p in pills
        )
        pills_html = f'<div class="cover-pills">{pill_items}</div>'
    title_b_html = (
        f'<h1 class="cover-title-b">{title_b}</h1>' if title_b else ""
    )
    subhead_html = (
        f'<p class="cover-subhead">{subhead}</p>' if subhead else ""
    )
    return (
        '<section class="pg pg-dark">'
        f'{brand_html}'
        f'{kicker_html}'
        f'<h1 class="cover-title-a">{title_a}</h1>'
        f'{title_b_html}'
        f'{subhead_html}'
        f'{pills_html}'
        '</section>'
    )


def _pdf_render_argument(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    subhead = _esc(sec.get("subhead"))
    section_heading = _esc(sec.get("section_heading"))
    section_intro = _esc(sec.get("section_intro"))
    bullets = [_esc(b) for b in (sec.get("bullets") or [])]
    takeaway = _esc(sec.get("takeaway"))
    kicker_html = f'<p class="kicker">{kicker}</p>' if kicker else ""
    subhead_html = (
        f'<p class="subhead">{subhead}</p>' if subhead else ""
    )
    head_html = (
        f'<h2 class="arg-section-heading">{section_heading}</h2>'
        if section_heading else ""
    )
    intro_html = (
        f'<p class="arg-intro">{section_intro}</p>'
        if section_intro else ""
    )
    bullets_html = "".join(f"<li>{b}</li>" for b in bullets)
    takeaway_html = (
        f'<div class="arg-takeaway">{takeaway}</div>'
        if takeaway else ""
    )
    return (
        '<section class="pg pg-light">'
        f'{kicker_html}'
        f'<h1 class="action-title">{title_text}</h1>'
        f'{subhead_html}'
        '<div class="arg-card">'
        f'{head_html}'
        f'{intro_html}'
        f'<ul class="arg-bullets">{bullets_html}</ul>'
        '</div>'
        f'{takeaway_html}'
        '</section>'
    )


def _pdf_render_stat_grid(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    subhead = _esc(sec.get("subhead"))
    stats = list(sec.get("stats") or [])
    detail_heading = _esc(sec.get("detail_heading"))
    detail_body = _esc(sec.get("detail_body"))
    accents = ("accent-blue", "accent-cyan", "accent-purple", "accent-green")
    cards = []
    for i in range(4):
        s = stats[i] if i < len(stats) and isinstance(stats[i], Mapping) else {}
        cards.append(
            f'<div class="stat-card {accents[i]}">'
            f'<p class="stat-value">{_esc(s.get("value"))}</p>'
            f'<p class="stat-label">{_esc(s.get("label"))}</p>'
            '</div>'
        )
    kicker_html = f'<p class="kicker">{kicker}</p>' if kicker else ""
    subhead_html = (
        f'<p class="subhead">{subhead}</p>' if subhead else ""
    )
    detail_html = ""
    if detail_heading or detail_body:
        detail_html = (
            '<div class="stat-detail">'
            + (
                f'<h2 class="stat-detail-heading">{detail_heading}</h2>'
                if detail_heading else ""
            )
            + (
                f'<p class="stat-detail-body">{detail_body}</p>'
                if detail_body else ""
            )
            + '</div>'
        )
    return (
        '<section class="pg pg-light">'
        f'{kicker_html}'
        f'<h1 class="action-title">{title_text}</h1>'
        f'{subhead_html}'
        f'<div class="stat-grid">{"".join(cards)}</div>'
        f'{detail_html}'
        '</section>'
    )


def _pdf_render_features_2x2(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    subhead = _esc(sec.get("subhead"))
    features = list(sec.get("features") or [])
    banner_text = _esc(sec.get("banner_text"))
    icon_classes = ("icon-blue", "icon-cyan", "icon-purple", "icon-green")
    glyph_defaults = ("✱", "◆", "▲", "●")
    cards = []
    for i in range(4):
        f = features[i] if i < len(features) and isinstance(features[i], Mapping) else {}
        glyph = _esc(f.get("icon") or glyph_defaults[i])
        cards.append(
            '<div class="feat-card">'
            '<div class="feat-head">'
            f'<span class="feat-icon {icon_classes[i]}">{glyph}</span>'
            f'<h3 class="feat-title">{_esc(f.get("title"))}</h3>'
            '</div>'
            f'<p class="feat-body">{_esc(f.get("body"))}</p>'
            '</div>'
        )
    kicker_html = (
        f'<p class="kicker kicker-purple">{kicker}</p>' if kicker else ""
    )
    subhead_html = (
        f'<p class="subhead">{subhead}</p>' if subhead else ""
    )
    banner_html = (
        f'<div class="feat-banner">{banner_text}</div>' if banner_text else ""
    )
    return (
        '<section class="pg pg-light">'
        f'{kicker_html}'
        f'<h1 class="action-title">{title_text}</h1>'
        f'{subhead_html}'
        f'<div class="feat-grid">{"".join(cards)}</div>'
        f'{banner_html}'
        '</section>'
    )


def _pdf_render_industries_2x2(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    subhead = _esc(sec.get("subhead"))
    industries = list(sec.get("industries") or [])
    rule_classes = ("rule-blue", "rule-green", "rule-orange", "rule-purple")
    icon_classes = ("icon-blue", "icon-green", "icon-orange", "icon-purple")
    glyph_defaults = ("◆", "▲", "●", "■")
    cards = []
    for i in range(4):
        ind = industries[i] if i < len(industries) and isinstance(industries[i], Mapping) else {}
        glyph = _esc(ind.get("icon") or glyph_defaults[i])
        uses_html = "".join(
            f"<li>{_esc(u)}</li>"
            for u in (ind.get("uses") or [])
        )
        cards.append(
            f'<div class="ind-card {rule_classes[i]}">'
            '<div class="ind-head">'
            f'<span class="ind-icon {icon_classes[i]}">{glyph}</span>'
            f'<h3 class="ind-name">{_esc(ind.get("name"))}</h3>'
            '</div>'
            f'<ul class="ind-uses">{uses_html}</ul>'
            '</div>'
        )
    kicker_html = f'<p class="kicker">{kicker}</p>' if kicker else ""
    subhead_html = (
        f'<p class="subhead">{subhead}</p>' if subhead else ""
    )
    return (
        '<section class="pg pg-light">'
        f'{kicker_html}'
        f'<h1 class="action-title">{title_text}</h1>'
        f'{subhead_html}'
        f'<div class="ind-grid">{"".join(cards)}</div>'
        '</section>'
    )


def _pdf_render_sovereignty_dark(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    subhead = _esc(sec.get("subhead"))
    cards = list(sec.get("cards") or [])
    gov_heading = _esc(sec.get("gov_heading"))
    governance = list(sec.get("governance") or [])
    card_html = []
    for i in range(4):
        c = cards[i] if i < len(cards) and isinstance(cards[i], Mapping) else {}
        card_html.append(
            '<div class="sov-card">'
            f'<h3 class="sov-card-title">{_esc(c.get("title"))}</h3>'
            f'<p class="sov-card-body">{_esc(c.get("body"))}</p>'
            '</div>'
        )
    gov_html = ""
    if gov_heading or governance:
        items = []
        for g in governance[:4]:
            if not isinstance(g, Mapping):
                continue
            items.append(
                '<div class="sov-gov-item">'
                '<span class="sov-gov-check">✓</span>'
                '<div>'
                f'<p class="sov-gov-title">{_esc(g.get("title"))}</p>'
                f'<p class="sov-gov-body">{_esc(g.get("body"))}</p>'
                '</div>'
                '</div>'
            )
        gov_html = (
            '<div class="sov-gov">'
            + (
                f'<p class="sov-gov-heading">{gov_heading}</p>'
                if gov_heading else ""
            )
            + f'<div class="sov-gov-grid">{"".join(items)}</div>'
            + '</div>'
        )
    kicker_html = (
        f'<p class="kicker kicker-cyan">{kicker}</p>' if kicker else ""
    )
    subhead_html = (
        f'<p class="subhead on-dark">{subhead}</p>' if subhead else ""
    )
    return (
        '<section class="pg pg-dark">'
        f'{kicker_html}'
        f'<h1 class="action-title on-dark">{title_text}</h1>'
        f'{subhead_html}'
        f'<div class="sov-grid">{"".join(card_html)}</div>'
        f'{gov_html}'
        '</section>'
    )


def _pdf_render_closing_dark(sec: Mapping[str, Any]) -> str:
    kicker = _esc(sec.get("kicker"))
    title_text = _esc(sec.get("title"))
    reasons = list(sec.get("reasons") or [])
    cta = _esc(sec.get("cta"))
    card_html = []
    for i in range(4):
        r = reasons[i] if i < len(reasons) and isinstance(reasons[i], Mapping) else None
        if r is None:
            continue
        num = _esc(r.get("number") or f"{i + 1:02d}")
        card_html.append(
            '<div class="close-card">'
            f'<div class="close-num">{num}</div>'
            '<div>'
            f'<h3 class="close-title">{_esc(r.get("title"))}</h3>'
            f'<p class="close-body">{_esc(r.get("body"))}</p>'
            '</div>'
            '</div>'
        )
    kicker_html = (
        f'<p class="kicker kicker-cyan">{kicker}</p>' if kicker else ""
    )
    cta_html = f'<div class="close-cta">{cta}</div>' if cta else ""
    return (
        '<section class="pg pg-dark">'
        f'{kicker_html}'
        f'<h1 class="action-title on-dark">{title_text}</h1>'
        f'<div class="close-grid">{"".join(card_html)}</div>'
        f'{cta_html}'
        '</section>'
    )


_PDF_EXEC_RENDERERS = {
    "cover":             _pdf_render_cover,
    "argument":          _pdf_render_argument,
    "stat_grid":         _pdf_render_stat_grid,
    "features_2x2":      _pdf_render_features_2x2,
    "industries_2x2":    _pdf_render_industries_2x2,
    "sovereignty_dark":  _pdf_render_sovereignty_dark,
    "closing_dark":      _pdf_render_closing_dark,
}


def make_pdf_exec(
    sections: list[Mapping[str, Any]],
    *,
    title: str = "report",
    summary: str = "",
    auto_publish: bool = True,
) -> Path:
    """Render a structured exec-PDF spec to a polished A4 PDF.

    ``sections`` is a list of section dicts; each dict's ``type`` selects
    one of the 7 executive page templates (mirrors the Citra-Service web
    pipeline's ``exec_pg_*`` family):

        cover, argument, stat_grid, features_2x2,
        industries_2x2, sovereignty_dark, closing_dark

    Example::

        sections = [
            {"type": "cover", "title_a": "FY26 plan", "title_b": "Three forces",
             "kicker": "Board pre-read", "pills": ["Growth","Margin","Cash"]},
            {"type": "argument",
             "kicker": "MID-MARKET",
             "title": "Consolidation is permanent",
             "bullets": ["Top-5 share up 11 quarters running", "..."]},
            {"type": "stat_grid",
             "title": "Four numbers describe FY26",
             "stats": [
                 {"value": "$4.2B", "label": "Revenue"},
                 {"value": "31%",   "label": "Op margin"},
                 {"value": "0.92",  "label": "Cash conversion"},
                 {"value": "11",    "label": "Quarters of growth"},
             ]},
            # ... 9 more sections, e.g. features_2x2, industries_2x2,
            # sovereignty_dark, more argument pages, closing_dark ...
        ]
        path = make_pdf_exec(sections, title="FY26 plan",
                             summary="Executive board pre-read")

    Each section becomes one A4 page. Pages flow in the order given.
    Image / chart fields are not part of this spec — the exec family is
    typography + colour-block only. Auto-publishes by default.

    See ``SKILL_REPORTS_EXEC.md`` for the full section catalogue, rhythm
    rules, and a 12-page worked example.
    """
    from weasyprint import HTML

    pages: list[str] = []
    for i, sec in enumerate(sections, start=1):
        if not isinstance(sec, Mapping):
            raise TypeError(f"section #{i} is not a dict")
        kind = str(sec.get("type") or "").lower()
        if kind not in _PDF_EXEC_RENDERERS:
            raise ValueError(
                f"section #{i}: unknown type {kind!r}; "
                f"choose one of {_PDF_EXEC_SECTIONS}"
            )
        pages.append(_PDF_EXEC_RENDERERS[kind](sec))

    doc = (
        '<!DOCTYPE html><html><head><meta charset="utf-8">'
        f'<title>{_esc(title)}</title>'
        f'<style>{_EXEC_PDF_CSS}</style>'
        f'</head><body>{"".join(pages)}</body></html>'
    )

    out = _stamped_path(title, "pdf")
    HTML(string=doc, base_url=str(_reports_dir())).write_pdf(str(out))
    _maybe_publish(out, kind="pdf", title=title, summary=summary,
                   auto_publish=auto_publish)
    return out


# ------------------------------- Excel -------------------------------------
def make_excel(
    sheets: Mapping[str, Iterable[Mapping[str, Any]]],
    *,
    title: str = "report",
    summary: str = "",
    auto_publish: bool = True,
) -> Path:
    """Write an .xlsx with one sheet per ``sheets`` entry.

    Each sheet's rows are dicts; the union of keys across the first 100 rows
    becomes the header. Auto-publishes by default.
    """
    from openpyxl import Workbook  # heavy import, defer

    wb = Workbook()
    # openpyxl creates a default sheet; remove it then re-add ours.
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=(name or "Sheet")[:31])
        rows_list = list(rows)
        headers: list[str] = []
        seen: set[str] = set()
        for row in rows_list[:100]:
            for k in row.keys():
                if k not in seen:
                    seen.add(k)
                    headers.append(k)
        if headers:
            ws.append(headers)
            for row in rows_list:
                ws.append([row.get(h) for h in headers])
        else:
            ws.append(["(empty)"])
    out = _stamped_path(title, "xlsx")
    wb.save(str(out))
    _maybe_publish(out, kind="xlsx", title=title, summary=summary, auto_publish=auto_publish)
    return out


# ------------------------------- CSV ---------------------------------------
def make_csv(
    rows: Iterable[Mapping[str, Any]],
    *,
    title: str = "report",
    summary: str = "",
    auto_publish: bool = True,
) -> Path:
    rows_list = [dict(r) for r in rows]
    headers: list[str] = []
    seen: set[str] = set()
    for row in rows_list[:100]:
        for k in row.keys():
            if k not in seen:
                seen.add(k)
                headers.append(k)
    out = _stamped_path(title, "csv")
    buf = io.StringIO()
    w = csv.DictWriter(buf, fieldnames=headers or ["value"])
    w.writeheader()
    for row in rows_list:
        w.writerow({h: row.get(h) for h in (headers or ["value"])})
    out.write_text(buf.getvalue(), encoding="utf-8")
    _maybe_publish(out, kind="csv", title=title, summary=summary, auto_publish=auto_publish)
    return out


# ------------------------------- PPTX --------------------------------------
# Corporate palette + type scale used by make_pptx. Kept conservative
# (single-accent navy + warm-grey neutrals) so a deck looks polished
# without theme work. Override via the ``theme`` arg if you need brand
# colors. Hex is converted to python-pptx RGBColor inline.
_DEFAULT_PPTX_THEME: dict[str, Any] = {
    "background":   "FFFFFF",   # slide background
    "text_primary": "1A1F2E",   # body / titles
    "text_muted":   "5E6877",   # subtitles, footers
    "accent":       "1E4D8C",   # title underline, dividers, callouts
    "rule":         "E5E8EE",   # 1px separator lines
    "font_head":    "Calibri",
    "font_body":    "Calibri",
    "size_title":   40,         # pt — title slide H1
    "size_h1":      28,         # pt — content slide title
    "size_h2":      18,         # pt — section divider title
    "size_body":    16,         # pt — content bullets / paragraph
    "size_kicker":  12,         # pt — section labels, footers
}

# Allowed layout types — keep this tight so the agent doesn't invent
# layouts the renderer can't honor.
_PPTX_LAYOUTS = (
    "freeform",              # no template — a list of positioned elements
    "title",                 # opener: title + subtitle + (optional) hero image
    "section_divider",       # full-bleed accent block + section name
    "content_bullets",       # H1 + bullet list (4–6 items max)
    "content_with_image",    # H1 + bullets (left, 55%) + image (right, 45%)
    "two_column",            # H1 + two body columns (50/50)
    "full_bleed_image",      # edge-to-edge image with overlay caption
    "quote",                 # large quote + attribution
    "closing",               # thank-you / next-steps / CTA + (optional) image
    # Executive corporate-boardroom templates — parity with the Citra-Service
    # web pipeline's exec_* family (slide_templates.py). All are typography +
    # colour-block only; any 'image' field on the slide spec is ignored.
    "exec_title_dark",       # dark cover — two-tone headline + pill labels
    "exec_three_pillars",    # 3 two-tone pillar cards + claim banner
    "exec_action_card",      # left 'how' card + right navy stat card
    "exec_argument",         # workhorse — kicker + action-title + bullets
    "exec_stat_grid_4",      # 4 KPI cards with accent bars
    "exec_features_2x2",     # 2×2 feature grid with icon circles
    "exec_industries_2x2",   # 2×2 industry grid with side rules
    "exec_sovereignty_dark", # dark — 4 trust cards + governance panel
    "exec_chat_example",     # left chat card + right 3 stat blocks
    "exec_closing_dark",     # dark — 2×2 numbered cards + cyan CTA banner
)

# Footer page numbers are suppressed on these layouts (each carries its
# own design language; a small footer would be noise).
_EXEC_LAYOUTS = frozenset({
    "exec_title_dark", "exec_three_pillars", "exec_action_card",
    "exec_argument", "exec_stat_grid_4", "exec_features_2x2",
    "exec_industries_2x2", "exec_sovereignty_dark", "exec_chat_example",
    "exec_closing_dark",
})

# Fixed palette for the executive family — these are NOT theme-driven
# (the web side bakes them into each template's slot definitions, and
# we mirror that 1:1 so decks look identical across surfaces).
_EXEC = {
    "navy_bg":     "0B1020",  # dark slide background
    "navy_card":   "111827",  # dark cards on navy bg
    "navy_border": "1E293B",
    "light_bg":    "F8FAFC",  # light slide background
    "white":       "FFFFFF",
    "ink":         "0F172A",  # primary text on light
    "slate":       "475569",  # secondary text on light
    "slate_dark":  "334155",
    "slate_soft":  "94A3B8",
    "slate_light": "CBD5E1",  # secondary text on dark
    "blue":        "2563EB",  # primary accent (left/odd cards)
    "cyan":        "06B6D4",
    "cyan_bright": "22D3EE",  # accent on dark bgs
    "purple":      "8B5CF6",
    "green":       "10B981",
    "orange":      "F59E0B",
    "blue_soft":   "BFDBFE",
    "cyan_soft":   "CFFAFE",
    "purple_soft": "DDD6FE",
    "green_soft":  "DCFCE7",
    "green_ink":   "065F46",
    "chat_panel":  "EEF2F7",
}


# ---- Freeform auto-layout: overlap resolution -----------------------------
# Mirrors the legacy Citra-Service presentation_api.fix_overlapping_elements
# pass. The agent positions freeform elements by hand, and an LLM routinely
# stacks two boxes at the same y or under-sizes a text box. This
# deterministic pass pushes vertically-colliding content elements apart so
# a deck never ships with overlapping text.

def _estimate_freeform_height(el: Mapping[str, Any]) -> float:
    """Best-effort rendered height of a freeform element, in 960×540 px.

    Text height is estimated from character wrapping (the agent's declared
    ``height`` is often too small for the content); other element types
    use their declared ``height``.
    """
    etype = str(el.get("type") or "text").lower()
    declared = float(el.get("height", 0) or 0)
    if etype == "text":
        text = el.get("text")
        if not text and isinstance(el.get("bullets"), (list, tuple)):
            text = "\n".join(str(b) for b in el["bullets"])
        text = str(text or "")
        if not text:
            return declared or 40.0
        fs = float(el.get("fontSize", 20) or 20)
        w = float(el.get("width", 400) or 400)
        cpl = max(1, int(w / (fs * 0.48)))            # chars per line
        lines = 0
        for ln in text.split("\n"):
            lines += max(1, -(-len(ln) // cpl))       # ceil division
        return max(lines * fs * 1.4, declared)
    return declared or 60.0


def _resolve_freeform_overlaps(
    elements: list[dict[str, Any]],
    *,
    canvas_h: float = 540.0,
    gap: float = 14.0,
    margin: float = 24.0,
) -> None:
    """Push vertically-overlapping content elements apart, in place.

    Skipped from the pass:

    - ``shape`` elements — backgrounds / accent bars / dividers.
    - **Full-bleed elements** of ANY type (an ``image``/``svg``/``shape``
      that covers ~the whole 960×540 canvas) — a full-bleed hero image is
      a background, not content. Including it would make every text box
      "overlap" it and cascade the whole slide to the bottom.

    Two elements only count as colliding when their X-ranges genuinely
    overlap (≥40% of the narrower box), so side-by-side cards and a
    left-text / right-chart split are left alone.
    """
    def _is_background(e: dict[str, Any]) -> bool:
        if str(e.get("type") or "text").lower() == "shape":
            return True
        # Near-full-canvas element → it is a backdrop, not content.
        return (
            float(e.get("x", 0) or 0) <= 16
            and float(e.get("y", 0) or 0) <= 16
            and float(e.get("width", 0) or 0) >= 900
            and float(e.get("height", 0) or 0) >= 500
        )

    items = [
        e for e in elements
        if isinstance(e, dict) and not _is_background(e)
    ]
    if len(items) < 2:
        return
    eff = {id(e): _estimate_freeform_height(e) for e in items}
    items.sort(key=lambda e: float(e.get("y", 0) or 0))
    for i in range(1, len(items)):
        curr = items[i]
        cy = float(curr.get("y", 0) or 0)
        cx = float(curr.get("x", 0) or 0)
        cw = float(curr.get("width", 0) or 0) or 1.0
        for j in range(i):
            prev = items[j]
            px_ = float(prev.get("x", 0) or 0)
            pw = float(prev.get("width", 0) or 0) or 1.0
            overlap_x = min(cx + cw, px_ + pw) - max(cx, px_)
            if overlap_x < 0.4 * min(cw, pw):
                continue
            p_bottom = float(prev.get("y", 0) or 0) + eff[id(prev)]
            if cy < p_bottom + gap:
                cy = p_bottom + gap
        # Keep the element on-canvas.
        cy = min(cy, max(margin, canvas_h - margin - eff[id(curr)]))
        curr["y"] = round(cy)


def make_pptx(
    slides: "list[Mapping[str, Any]] | str | Path",
    *,
    title: str = "presentation",
    summary: str = "",
    theme: Mapping[str, Any] | None = None,
    auto_publish: bool = True,
) -> Path:
    """Render a slide spec to a .pptx file.

    ``slides`` may be a list of slide dicts **or** the path to a JSON file
    holding the deck. Authoring the deck as a JSON file (then
    ``make_pptx("deck.json")``) is strongly preferred over hand-writing a
    nested Python literal — the file is either a bare slides list or
    ``{"title": ..., "summary": ..., "slides": [...]}``.

    ``slides`` is a list of dicts; each dict is one slide.

    **Freeform (recommended).** A slide with an ``elements`` list is a
    blank 960×540 canvas you compose yourself — no template::

        {"background": "#0B1020",
         "elements": [
           {"type": "text",  "text": "Q4 Outlook", "x": 80, "y": 200,
            "width": 800, "height": 110, "fontSize": 54, "bold": True,
            "color": "#FFFFFF", "align": "left"},
           {"type": "shape", "shape": "rect", "x": 80, "y": 320,
            "width": 90, "height": 8, "fill": "#06B6D4"},
           {"type": "image", "x": 540, "y": 70, "width": 360, "height": 400,
            "image": <bytes | Path | data-URI | "<svg>…</svg>">},
           {"type": "chart", "x": 80, "y": 350, "width": 420, "height": 160,
            "chart": {"type": "bar", "data": {...}}},  # a Chart.js config
           {"type": "card", "x": 520, "y": 350, "width": 360, "height": 150,
            "fill": "#111827", "accent": "#06B6D4", "icon": "📈",
            "title": "Coverage", "body": "Up 4pp on GTM discipline"},
         ]}

    Element types — all positioned in the 960×540 design space, with
    ``zIndex`` controlling paint order:

    - ``text``  — ``text`` or ``bullets[]``; ``fontSize``, ``color``,
      ``bold``, ``align``, ``valign``, ``fontFamily``.
    - ``image`` — ``image`` (bytes / ``Path`` / ``data:`` URI / ``<svg>``).
    - ``svg``   — ``svg`` (raw ``<svg>`` markup) — a simple drawing/diagram.
    - ``chart`` — ``chart`` = a Chart.js config, rendered via the baked-in
      renderer.
    - ``card``  — rounded panel: ``fill``, ``accent``, ``icon``, ``title``,
      ``body`` (or ``bullets[]``); ``lineColor``, ``titleColor``,
      ``bodyColor``.
    - ``shape`` — ``shape`` (rect/rounded/ellipse/line), ``fill``,
      ``lineColor``, ``lineWidth``, optional ``text``.

    **Templated (legacy).** A slide with a ``layout`` key uses a fixed
    template — see ``_PPTX_LAYOUTS`` for the names::

        {"layout": "title",
         "title": "Q4 Market Outlook",
         "subtitle": "Where the next $100M of growth comes from",
         "image": <bytes | citra_toolkit.image.GeneratedImage | None>}

        {"layout": "content_bullets",
         "title": "Three forces reshaping demand",
         "bullets": ["...", "...", "..."]}

        {"layout": "content_with_image",
         "title": "Mid-market is consolidating fast",
         "bullets": ["...", "..."],
         "image": <bytes>,
         "caption": "AI-generated illustration"}

        {"layout": "section_divider",
         "title": "Recommendations",
         "kicker": "Part 3 of 4"}

        {"layout": "two_column",
         "title": "Buy vs build",
         "left":  {"heading": "Buy",  "body": "..."},
         "right": {"heading": "Build","body": "..."}}

        {"layout": "full_bleed_image",
         "image": <bytes>,
         "caption": "Boardroom in mid-discussion"}

        {"layout": "quote",
         "quote": "The market punishes the second-best for being on time.",
         "attribution": "— A. Karpathy"}

        {"layout": "closing",
         "title": "What we need from you",
         "bullets": ["Sign off on the FY27 plan", "Approve the hire"],
         "image": <bytes | None>}

    Image fields accept raw ``bytes``, a ``citra_toolkit.image.GeneratedImage``
    (the ``.data`` is pulled automatically), or ``None``. Generate
    images with ``citra_toolkit.image.generate(...)`` and pass through.

    Auto-publishes by default (surfaces as a downloadable .pptx in the
    chat). Pass ``auto_publish=False`` to defer.
    """
    # Accept a path to a JSON deck spec — the agent authors deck.json with
    # the file tool (plain JSON, no Python-literal escaping to corrupt) and
    # renders it in one line. The file is a slides list, or a dict with
    # `slides` plus optional `title` / `summary`.
    if isinstance(slides, (str, Path)):
        import json as _json
        with open(slides, encoding="utf-8") as _f:
            _spec = _json.load(_f)
        if isinstance(_spec, Mapping):
            if title == "presentation" and _spec.get("title"):
                title = str(_spec["title"])
            if not summary and _spec.get("summary"):
                summary = str(_spec["summary"])
            slides = list(_spec.get("slides") or [])
        elif isinstance(_spec, list):
            slides = _spec
        else:
            raise TypeError(
                "deck JSON must be a slides list or a {slides, title?, summary?} object"
            )

    from pptx import Presentation  # type: ignore[import-not-found]
    from pptx.util import Inches, Pt, Emu  # type: ignore[import-not-found]
    from pptx.dml.color import RGBColor  # type: ignore[import-not-found]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-not-found]
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import-not-found]

    t: dict[str, Any] = {**_DEFAULT_PPTX_THEME, **(theme or {})}

    def hex_to_rgb(h: str) -> "RGBColor":
        h = str(h).strip().lstrip("#")
        # Expand 3-digit shorthand (#fff -> #ffffff) — agents and the web
        # canvas both emit it freely.
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except (ValueError, IndexError):
            return RGBColor(0x1A, 0x1A, 0x1A)  # safe near-black fallback

    def coerce_image(value: Any) -> bytes | None:
        if value is None:
            return None
        # Path to a rendered file — e.g. the return of make_chartjs_png /
        # make_chart_png, or any image the agent saved to disk.
        if isinstance(value, Path):
            try:
                return value.read_bytes()
            except Exception:  # noqa: BLE001
                return None
        if isinstance(value, str):
            # SVG markup — rasterise to PNG so python-pptx can embed it.
            if _looks_like_svg(value):
                try:
                    return svg_to_png(value)
                except Exception:  # noqa: BLE001
                    return None
            # data: URI — e.g. the `image_data` field of the
            # citra_image_generate MCP tool's response.
            if value.startswith("data:"):
                try:
                    _, _, payload = value.partition(",")
                    return base64.b64decode(payload)
                except Exception:  # noqa: BLE001
                    return None
            # Bare filesystem path string.
            try:
                p = Path(value)
                if p.is_file():
                    return p.read_bytes()
            except Exception:  # noqa: BLE001
                pass
            return None
        if isinstance(value, (bytes, bytearray)):
            if _looks_like_svg(value):
                try:
                    return svg_to_png(bytes(value))
                except Exception:  # noqa: BLE001
                    return None
            return bytes(value)
        # Duck-type GeneratedImage / anything with a .data bytes attribute.
        data = getattr(value, "data", None)
        if isinstance(data, (bytes, bytearray)):
            return bytes(data)
        return None

    def set_text(frame, text: str, *, size: int, color_hex: str,
                 bold: bool = False, align: int | None = None,
                 font: str | None = None) -> None:
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font or t["font_body"]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color_hex)

    def add_bullets(frame, bullets: list[str], *, size: int,
                    color_hex: str) -> None:
        frame.clear()
        frame.word_wrap = True
        for i, b in enumerate(bullets):
            p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            run = p.add_run()
            run.text = b
            run.font.name = t["font_body"]
            run.font.size = Pt(size)
            run.font.color.rgb = hex_to_rgb(color_hex)
            # Space after for readability.
            p.space_after = Pt(8)

    def paint_background(slide, color_hex: str) -> None:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.line.fill.background()
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(color_hex)
        # Send to back: pptx has no z-order API, so we re-order via XML.
        spTree = bg._element.getparent()  # noqa: SLF001
        spTree.remove(bg._element)        # noqa: SLF001
        spTree.insert(2, bg._element)     # noqa: SLF001

    def add_accent_rule(slide, *, top: Emu, left: Emu, width: Emu,
                        height: Emu) -> None:
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        rect.line.fill.background()
        rect.fill.solid()
        rect.fill.fore_color.rgb = hex_to_rgb(t["accent"])

    # ---- Helpers shared by exec_* layouts ---------------------------------
    # The exec family uses absolute positioning in a 960×540 design canvas
    # (matching slide_templates.py on the web side). 16:9 means a single
    # scale factor maps both axes: 13.333 inches / 960 px = 7.5 / 540.
    _PX = 13.333 / 960  # inches per template-pixel

    def px(n: float) -> Emu:
        return Inches(n * _PX)

    def add_rect(slide, *, x: float, y: float, w: float, h: float,
                 fill_hex: str, stroke_hex: str | None = None,
                 stroke_w: float = 0.75):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h)
        )
        if stroke_hex is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = hex_to_rgb(stroke_hex)
            shape.line.width = Pt(stroke_w)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        return shape

    def add_rounded_rect(slide, *, x: float, y: float, w: float, h: float,
                         fill_hex: str, stroke_hex: str | None = None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
        )
        if stroke_hex is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = hex_to_rgb(stroke_hex)
            shape.line.width = Pt(0.75)
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        return shape

    def add_circle(slide, *, x: float, y: float, size: float,
                   fill_hex: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, px(x), px(y), px(size), px(size)
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
        return shape

    def text_at(slide, *, x: float, y: float, w: float, h: float,
                text: str, size: int, color_hex: str,
                bold: bool = False, align: str = "left",
                font: str | None = None,
                anchor: str = "top"):
        tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(anchor, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.name = font or t["font_body"]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color_hex)
        return tb

    def bullets_at(slide, *, x: float, y: float, w: float, h: float,
                   bullets: list, size: int, color_hex: str,
                   bullet_color_hex: str | None = None,
                   lh: float = 1.55, style: str = "dot"):
        tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        bc = bullet_color_hex or color_hex
        prefix = {"dot": "• ", "check": "✓ ", "dash": "— "}.get(style, "• ")
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            rb = p.add_run()
            rb.text = prefix
            rb.font.name = t["font_body"]
            rb.font.size = Pt(size)
            rb.font.bold = True
            rb.font.color.rgb = hex_to_rgb(bc)
            rt = p.add_run()
            rt.text = str(b)
            rt.font.name = t["font_body"]
            rt.font.size = Pt(size)
            rt.font.color.rgb = hex_to_rgb(color_hex)
            p.space_after = Pt(max(2, int(size * (lh - 1))))
        return tb

    # ---- Freeform layout --------------------------------------------------
    # No template: the slide is a list of absolutely-positioned elements in
    # the same 960×540 design canvas the exec_* family uses (mirrors the
    # Citra-Service web presentation builder's PRESENTATION_CANVAS). The
    # agent designs every slide itself — text boxes, images, charts, shapes.
    _ALIGN = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "centre": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    _VALIGN = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
    }
    _SHAPES = {
        "rect": MSO_SHAPE.RECTANGLE, "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL, "oval": MSO_SHAPE.OVAL,
        "circle": MSO_SHAPE.OVAL,
    }

    def render_freeform(slide, elements: Any) -> None:
        if not isinstance(elements, (list, tuple)):
            return
        # Shallow-copy each element so the caller's slide spec isn't
        # mutated, then push apart any vertically-overlapping content.
        els = [dict(e) for e in elements if isinstance(e, Mapping)]
        _resolve_freeform_overlaps(els)
        # zIndex controls paint order (higher = on top); default to the
        # element's position in the list.
        ordered = sorted(els, key=lambda e: e.get("zIndex", 0))
        for el in ordered:
            etype = str(el.get("type") or "text").lower()
            x = float(el.get("x", 0) or 0)
            y = float(el.get("y", 0) or 0)
            w = float(el.get("width", 200) or 0)
            h = float(el.get("height", 80) or 0)

            if etype == "text":
                box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
                tf = box.text_frame
                va = _VALIGN.get(str(el.get("valign") or "top").lower())
                if va is not None:
                    tf.vertical_anchor = va
                bullets = el.get("bullets")
                if isinstance(bullets, (list, tuple)) and bullets:
                    add_bullets(
                        tf, [str(b) for b in bullets],
                        size=int(el.get("fontSize", 18) or 18),
                        color_hex=str(el.get("color") or t["text_primary"]),
                    )
                else:
                    set_text(
                        tf, str(el.get("text") or ""),
                        size=int(el.get("fontSize", 20) or 20),
                        color_hex=str(el.get("color") or t["text_primary"]),
                        bold=bool(el.get("bold")),
                        align=_ALIGN.get(str(el.get("align") or "left").lower()),
                        font=el.get("fontFamily"),
                    )

            elif etype in ("image", "svg", "svg_diagram"):
                src = el.get("image")
                if src is None:
                    src = el.get("src")
                if src is None:
                    src = el.get("svg") or el.get("svgContent")
                data = coerce_image(src)
                if data:
                    slide.shapes.add_picture(
                        BytesIO(data), px(x), px(y), width=px(w), height=px(h)
                    )

            elif etype == "chart":
                cfg = el.get("chart") or el.get("chartConfig")
                if isinstance(cfg, Mapping):
                    try:
                        png = _chartjs_render(
                            cfg,
                            width=max(200, int(w * 2)),
                            height=max(150, int(h * 2)),
                            background=str(el.get("background") or "white"),
                        )
                        slide.shapes.add_picture(
                            BytesIO(png), px(x), px(y), width=px(w), height=px(h)
                        )
                    except Exception:  # noqa: BLE001
                        # A chart that fails to render must not sink the deck.
                        pass

            elif etype == "card":
                # A card = rounded panel + optional accent bar / icon /
                # title / body, composed for the agent so it doesn't have
                # to position 4 sub-elements by hand.
                pad = 24.0
                fill = str(el.get("fill") or el.get("backgroundColor")
                           or t.get("surface") or "#F1F5F9")
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = hex_to_rgb(fill)
                border = (el.get("lineColor") or el.get("borderColor")
                          or el.get("stroke"))
                if border:
                    card.line.color.rgb = hex_to_rgb(str(border))
                    card.line.width = Pt(float(el.get("lineWidth", 1) or 1))
                else:
                    card.line.fill.background()
                accent = el.get("accent") or el.get("accentColor")
                if accent:
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(6)
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = hex_to_rgb(str(accent))
                    bar.line.fill.background()
                inner_x = x + pad
                inner_w = max(40.0, w - 2 * pad)
                cur_y = y + pad
                icon = el.get("icon")
                if icon:
                    # Emoji / short glyph or an <svg> string. Lucide-by-name
                    # isn't available in pptx — use an emoji or an `svg`/
                    # `image` element for a real icon.
                    ico = coerce_image(icon) if _looks_like_svg(icon) else None
                    if ico:
                        slide.shapes.add_picture(
                            BytesIO(ico), px(inner_x), px(cur_y),
                            width=px(40), height=px(40),
                        )
                    else:
                        ib = slide.shapes.add_textbox(
                            px(inner_x), px(cur_y), px(inner_w), px(40)
                        )
                        set_text(ib.text_frame, str(icon), size=26,
                                 color_hex=str(accent or el.get("color")
                                                or t["text_primary"]))
                    cur_y += 48
                title = el.get("title")
                if title:
                    tb = slide.shapes.add_textbox(
                        px(inner_x), px(cur_y), px(inner_w), px(34)
                    )
                    set_text(tb.text_frame, str(title),
                             size=int(el.get("titleSize", 18) or 18), bold=True,
                             color_hex=str(el.get("titleColor") or el.get("color")
                                            or t["text_primary"]))
                    cur_y += 38
                body = el.get("body") or el.get("description") or el.get("text")
                bullets = el.get("bullets")
                body_h = max(20.0, (y + h - pad) - cur_y)
                body_color = str(el.get("bodyColor") or t.get("text_muted")
                                 or "#475569")
                if isinstance(bullets, (list, tuple)) and bullets:
                    bx = slide.shapes.add_textbox(
                        px(inner_x), px(cur_y), px(inner_w), px(body_h)
                    )
                    add_bullets(bx.text_frame, [str(b) for b in bullets],
                                size=int(el.get("bodySize", 13) or 13),
                                color_hex=body_color)
                elif body:
                    bx = slide.shapes.add_textbox(
                        px(inner_x), px(cur_y), px(inner_w), px(body_h)
                    )
                    set_text(bx.text_frame, str(body),
                             size=int(el.get("bodySize", 13) or 13),
                             color_hex=body_color)

            elif etype == "shape":
                mso = _SHAPES.get(str(el.get("shape") or "rect").lower(),
                                  MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(mso, px(x), px(y), px(w), px(h))
                fill = el.get("fill")
                if fill:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb(str(fill))
                else:
                    shape.fill.background()
                stroke = el.get("lineColor") or el.get("stroke")
                if stroke:
                    shape.line.color.rgb = hex_to_rgb(str(stroke))
                    shape.line.width = Pt(float(el.get("lineWidth", 1) or 1))
                else:
                    shape.line.fill.background()
                txt = el.get("text")
                if txt:
                    set_text(
                        shape.text_frame, str(txt),
                        size=int(el.get("fontSize", 18) or 18),
                        color_hex=str(el.get("color") or t["text_primary"]),
                        bold=bool(el.get("bold")),
                        align=_ALIGN.get(str(el.get("align") or "center").lower()),
                    )

    prs = Presentation()
    # 16:9 — modern board-deck default.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    MARGIN = Inches(0.6)
    blank_layout = prs.slide_layouts[6]  # truly blank

    for idx, raw in enumerate(slides, start=1):
        if not isinstance(raw, Mapping):
            raise TypeError(f"slide #{idx} is not a dict")
        layout = str(raw.get("layout") or "").lower()

        # Freeform: the slide carries its own positioned elements. Detected
        # by an explicit layout, or simply by the presence of `elements`.
        if layout == "freeform" or (not layout and "elements" in raw):
            slide = prs.slides.add_slide(blank_layout)
            paint_background(slide, str(raw.get("background") or t["background"]))
            render_freeform(slide, raw.get("elements") or [])
            continue

        if not layout:
            layout = "content_bullets"
        if layout not in _PPTX_LAYOUTS:
            raise ValueError(
                f"slide #{idx}: unknown layout {layout!r}; "
                f"choose one of {_PPTX_LAYOUTS}"
            )

        slide = prs.slides.add_slide(blank_layout)
        paint_background(slide, t["background"])

        if layout == "title":
            # H1 centered vertically; accent rule beneath.
            title_text = str(raw.get("title") or "")
            subtitle = str(raw.get("subtitle") or "")
            image_bytes = coerce_image(raw.get("image"))

            if image_bytes:
                # Hero image fills right 45% of the slide.
                img_w = Inches(5.5)
                img_h = Inches(6.3)
                slide.shapes.add_picture(
                    BytesIO(image_bytes),
                    SW - img_w - MARGIN, Inches(0.6),
                    width=img_w, height=img_h,
                )
                text_right = SW - img_w - MARGIN * 2
            else:
                text_right = SW - MARGIN

            tx_w = text_right - MARGIN
            title_box = slide.shapes.add_textbox(
                MARGIN, Inches(2.6), tx_w, Inches(1.6)
            )
            set_text(title_box.text_frame, title_text,
                     size=t["size_title"], color_hex=t["text_primary"],
                     bold=True, font=t["font_head"])
            add_accent_rule(
                slide,
                top=Inches(4.45), left=MARGIN,
                width=Inches(0.9), height=Inches(0.07),
            )
            if subtitle:
                sub_box = slide.shapes.add_textbox(
                    MARGIN, Inches(4.7), tx_w, Inches(1.4)
                )
                set_text(sub_box.text_frame, subtitle,
                         size=t["size_body"], color_hex=t["text_muted"])

        elif layout == "section_divider":
            # Full-bleed accent block, large title, optional kicker.
            block = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, SW, SH
            )
            block.line.fill.background()
            block.fill.solid()
            block.fill.fore_color.rgb = hex_to_rgb(t["accent"])
            kicker = str(raw.get("kicker") or "")
            title_text = str(raw.get("title") or "")
            if kicker:
                kb = slide.shapes.add_textbox(
                    MARGIN, Inches(3.0), SW - MARGIN * 2, Inches(0.6)
                )
                set_text(kb.text_frame, kicker.upper(),
                         size=t["size_kicker"], color_hex="FFFFFF",
                         bold=True)
            tb = slide.shapes.add_textbox(
                MARGIN, Inches(3.4), SW - MARGIN * 2, Inches(1.8)
            )
            set_text(tb.text_frame, title_text,
                     size=t["size_title"], color_hex="FFFFFF",
                     bold=True, font=t["font_head"])

        elif layout == "content_bullets":
            title_text = str(raw.get("title") or "")
            bullets = list(raw.get("bullets") or [])
            tb = slide.shapes.add_textbox(
                MARGIN, MARGIN, SW - MARGIN * 2, Inches(0.9)
            )
            set_text(tb.text_frame, title_text,
                     size=t["size_h1"], color_hex=t["text_primary"],
                     bold=True, font=t["font_head"])
            add_accent_rule(
                slide,
                top=Inches(1.45), left=MARGIN,
                width=Inches(0.6), height=Inches(0.06),
            )
            body = slide.shapes.add_textbox(
                MARGIN, Inches(1.9), SW - MARGIN * 2, SH - Inches(2.6)
            )
            add_bullets(body.text_frame, bullets,
                        size=t["size_body"], color_hex=t["text_primary"])

        elif layout == "content_with_image":
            title_text = str(raw.get("title") or "")
            bullets = list(raw.get("bullets") or [])
            image_bytes = coerce_image(raw.get("image"))
            caption = str(raw.get("caption") or "")

            tb = slide.shapes.add_textbox(
                MARGIN, MARGIN, SW - MARGIN * 2, Inches(0.9)
            )
            set_text(tb.text_frame, title_text,
                     size=t["size_h1"], color_hex=t["text_primary"],
                     bold=True, font=t["font_head"])
            add_accent_rule(
                slide,
                top=Inches(1.45), left=MARGIN,
                width=Inches(0.6), height=Inches(0.06),
            )

            # Left column: bullets (55%)
            col_w = Inches(6.4)
            body = slide.shapes.add_textbox(
                MARGIN, Inches(1.9), col_w, SH - Inches(2.6)
            )
            add_bullets(body.text_frame, bullets,
                        size=t["size_body"], color_hex=t["text_primary"])

            # Right column: image (45%)
            img_x = MARGIN + col_w + Inches(0.4)
            img_w = SW - img_x - MARGIN
            img_h = Inches(5.1)
            if image_bytes:
                slide.shapes.add_picture(
                    BytesIO(image_bytes), img_x, Inches(1.9),
                    width=img_w, height=img_h,
                )
                if caption:
                    cap = slide.shapes.add_textbox(
                        img_x, Inches(1.9) + img_h + Inches(0.1),
                        img_w, Inches(0.4),
                    )
                    set_text(cap.text_frame, caption,
                             size=t["size_kicker"],
                             color_hex=t["text_muted"])

        elif layout == "two_column":
            title_text = str(raw.get("title") or "")
            left = raw.get("left") or {}
            right = raw.get("right") or {}
            tb = slide.shapes.add_textbox(
                MARGIN, MARGIN, SW - MARGIN * 2, Inches(0.9)
            )
            set_text(tb.text_frame, title_text,
                     size=t["size_h1"], color_hex=t["text_primary"],
                     bold=True, font=t["font_head"])
            add_accent_rule(
                slide,
                top=Inches(1.45), left=MARGIN,
                width=Inches(0.6), height=Inches(0.06),
            )
            col_w = (SW - MARGIN * 2 - Inches(0.5)) // 2
            for col_idx, col in enumerate((left, right)):
                x = MARGIN + col_idx * (col_w + Inches(0.5))
                heading = str(col.get("heading") or "")
                body_text = str(col.get("body") or "")
                head_box = slide.shapes.add_textbox(
                    x, Inches(1.9), col_w, Inches(0.6)
                )
                set_text(head_box.text_frame, heading,
                         size=t["size_h2"], color_hex=t["accent"],
                         bold=True)
                body_box = slide.shapes.add_textbox(
                    x, Inches(2.6), col_w, SH - Inches(3.2)
                )
                set_text(body_box.text_frame, body_text,
                         size=t["size_body"], color_hex=t["text_primary"])

        elif layout == "full_bleed_image":
            image_bytes = coerce_image(raw.get("image"))
            caption = str(raw.get("caption") or "")
            if not image_bytes:
                raise ValueError(
                    f"slide #{idx}: full_bleed_image requires an 'image'"
                )
            slide.shapes.add_picture(
                BytesIO(image_bytes), 0, 0, width=SW, height=SH,
            )
            if caption:
                # Translucent caption strip at the bottom.
                strip_h = Inches(0.9)
                strip = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, SH - strip_h, SW, strip_h,
                )
                strip.line.fill.background()
                strip.fill.solid()
                strip.fill.fore_color.rgb = hex_to_rgb(t["text_primary"])
                cap = slide.shapes.add_textbox(
                    MARGIN, SH - strip_h + Inches(0.2),
                    SW - MARGIN * 2, Inches(0.5),
                )
                set_text(cap.text_frame, caption,
                         size=t["size_body"], color_hex="FFFFFF")

        elif layout == "quote":
            quote = str(raw.get("quote") or "")
            attribution = str(raw.get("attribution") or "")
            qb = slide.shapes.add_textbox(
                MARGIN, Inches(2.2), SW - MARGIN * 2, Inches(3.2)
            )
            set_text(qb.text_frame, f"“{quote}”",
                     size=t["size_title"] - 6, color_hex=t["text_primary"],
                     bold=False, font=t["font_head"])
            if attribution:
                ab = slide.shapes.add_textbox(
                    MARGIN, Inches(5.5), SW - MARGIN * 2, Inches(0.6)
                )
                set_text(ab.text_frame, attribution,
                         size=t["size_body"], color_hex=t["text_muted"])
            add_accent_rule(
                slide,
                top=Inches(2.0), left=MARGIN,
                width=Inches(0.6), height=Inches(0.06),
            )

        elif layout == "closing":
            title_text = str(raw.get("title") or "Thank you")
            bullets = list(raw.get("bullets") or [])
            image_bytes = coerce_image(raw.get("image"))

            if image_bytes:
                # Image fills right 40%, text on left 60%.
                img_w = Inches(5.0)
                img_h = Inches(6.3)
                slide.shapes.add_picture(
                    BytesIO(image_bytes),
                    SW - img_w - MARGIN, Inches(0.6),
                    width=img_w, height=img_h,
                )
                text_right = SW - img_w - MARGIN * 2
            else:
                text_right = SW - MARGIN

            tx_w = text_right - MARGIN
            tb = slide.shapes.add_textbox(
                MARGIN, Inches(1.4), tx_w, Inches(1.2)
            )
            set_text(tb.text_frame, title_text,
                     size=t["size_title"] - 4, color_hex=t["text_primary"],
                     bold=True, font=t["font_head"])
            add_accent_rule(
                slide,
                top=Inches(2.85), left=MARGIN,
                width=Inches(0.9), height=Inches(0.07),
            )
            if bullets:
                bb = slide.shapes.add_textbox(
                    MARGIN, Inches(3.2), tx_w, SH - Inches(4.0)
                )
                add_bullets(bb.text_frame, bullets,
                            size=t["size_body"],
                            color_hex=t["text_primary"])

        # ====== Executive corporate-boardroom layouts =====================
        # All exec_* layouts use absolute coordinates from a 960×540 design
        # canvas (mirrors slide_templates.py on the web side). Palettes are
        # baked into each branch via the module-level _EXEC dict; theme arg
        # is intentionally ignored here. Image fields are also ignored —
        # exec slides are typography + colour-block only.

        elif layout == "exec_title_dark":
            paint_background(slide, _EXEC["navy_bg"])
            brand_chip = str(raw.get("brand_chip") or "").strip()
            kicker = str(raw.get("kicker") or "").strip()
            title_a = str(raw.get("title_a") or raw.get("title") or "")
            title_b = str(raw.get("title_b") or raw.get("subtitle") or "")
            subhead = str(raw.get("subhead") or "").strip()
            pills = list(raw.get("pills") or [])  # up to 3 short strings
            if brand_chip:
                # Faint stroke chip top-left.
                add_rect(slide, x=48, y=56, w=170, h=38,
                         fill_hex=_EXEC["navy_bg"],
                         stroke_hex=_EXEC["slate"], stroke_w=0.5)
                text_at(slide, x=48, y=56, w=170, h=38,
                        text=brand_chip.upper(), size=12,
                        color_hex=_EXEC["white"], bold=True,
                        align="center", anchor="middle")
            if kicker:
                text_at(slide, x=48, y=200, w=600, h=22,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["cyan_bright"], bold=True)
            text_at(slide, x=48, y=232, w=860, h=96,
                    text=title_a, size=60, color_hex=_EXEC["white"],
                    bold=True, font=t["font_head"])
            if title_b:
                text_at(slide, x=48, y=332, w=860, h=96,
                        text=title_b, size=60,
                        color_hex=_EXEC["cyan_bright"],
                        bold=True, font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=438, w=760, h=50,
                        text=subhead, size=15,
                        color_hex=_EXEC["slate_light"])
            for i, pill in enumerate(pills[:3]):
                px_left = (48, 350, 652)[i]
                add_rounded_rect(slide, x=px_left, y=500, w=260, h=28,
                                 fill_hex=_EXEC["navy_border"])
                text_at(slide, x=px_left, y=500, w=260, h=28,
                        text=str(pill), size=13,
                        color_hex=_EXEC["white"], bold=True,
                        align="center", anchor="middle")

        elif layout == "exec_three_pillars":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            pillars = list(raw.get("pillars") or [])
            banner = str(raw.get("banner") or "").strip()
            if kicker:
                text_at(slide, x=48, y=36, w=600, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["blue"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=90, text=title_text,
                    size=38, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=156, w=880, h=36,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            top_colors = (_EXEC["blue"], _EXEC["cyan"], _EXEC["purple"])
            label_colors = (_EXEC["blue_soft"], _EXEC["cyan_soft"],
                            _EXEC["purple_soft"])
            xs = (48, 340, 632)
            for i in range(3):
                pill = pillars[i] if i < len(pillars) else {}
                if not isinstance(pill, Mapping):
                    pill = {}
                add_rounded_rect(slide, x=xs[i], y=210, w=280, h=110,
                                 fill_hex=top_colors[i])
                if pill.get("label"):
                    text_at(slide, x=xs[i] + 22, y=230, w=240, h=18,
                            text=str(pill["label"]).upper(), size=11,
                            color_hex=label_colors[i], bold=True)
                if pill.get("title"):
                    text_at(slide, x=xs[i] + 22, y=254, w=240, h=56,
                            text=str(pill["title"]), size=17,
                            color_hex=_EXEC["white"], bold=True)
                add_rounded_rect(slide, x=xs[i], y=330, w=280, h=150,
                                 fill_hex=_EXEC["white"])
                if pill.get("body"):
                    text_at(slide, x=xs[i] + 18, y=348, w=244, h=120,
                            text=str(pill["body"]), size=13,
                            color_hex=_EXEC["slate_dark"])
            if banner:
                text_at(slide, x=48, y=498, w=864, h=30,
                        text=banner.upper(), size=13,
                        color_hex=_EXEC["blue"], bold=True,
                        align="center", anchor="middle")

        elif layout == "exec_action_card":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            left = raw.get("left") or {}
            right = raw.get("right") or {}
            if not isinstance(left, Mapping): left = {}
            if not isinstance(right, Mapping): right = {}
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["blue"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=60, text=title_text,
                    size=34, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=128, w=880, h=30,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            # Left "how it works" card.
            add_rounded_rect(slide, x=48, y=180, w=430, h=320,
                             fill_hex=_EXEC["white"])
            if left.get("title"):
                text_at(slide, x=72, y=204, w=380, h=28,
                        text=str(left["title"]), size=17,
                        color_hex=_EXEC["ink"], bold=True)
            steps = list(left.get("steps") or left.get("bullets") or [])
            for i, step in enumerate(steps[:6]):
                y_step = 244 + i * 36
                add_circle(slide, x=72, y=y_step, size=22,
                           fill_hex=_EXEC["blue"])
                text_at(slide, x=72, y=y_step, w=22, h=22,
                        text=str(i + 1), size=11,
                        color_hex=_EXEC["white"], bold=True,
                        align="center", anchor="middle")
                text_at(slide, x=102, y=y_step + 2, w=346, h=30,
                        text=str(step), size=13,
                        color_hex=_EXEC["slate_dark"])
            # Right navy stat card.
            add_rounded_rect(slide, x=502, y=180, w=410, h=320,
                             fill_hex=_EXEC["navy_bg"])
            if right.get("kicker"):
                text_at(slide, x=526, y=204, w=360, h=18,
                        text=str(right["kicker"]).upper(), size=11,
                        color_hex=_EXEC["cyan_bright"], bold=True)
            before = right.get("before") or {}
            after = right.get("after") or {}
            if not isinstance(before, Mapping): before = {}
            if not isinstance(after, Mapping): after = {}
            for col_idx, (col, value_color) in enumerate(
                ((before, _EXEC["white"]), (after, _EXEC["cyan_bright"]))
            ):
                cx = 526 + col_idx * 176
                if col.get("label"):
                    text_at(slide, x=cx, y=234, w=160, h=18,
                            text=str(col["label"]), size=13,
                            color_hex=_EXEC["cyan_bright"])
                if col.get("value") is not None:
                    text_at(slide, x=cx, y=256, w=160, h=56,
                            text=str(col["value"]), size=44,
                            color_hex=value_color, bold=True,
                            font=t["font_head"])
                if col.get("unit"):
                    text_at(slide, x=cx, y=316, w=160, h=18,
                            text=str(col["unit"]), size=13,
                            color_hex=_EXEC["slate_soft"])
            add_rect(slide, x=526, y=348, w=360, h=1,
                     fill_hex=_EXEC["navy_border"])
            list_label = str(right.get("list_label") or "").strip()
            if list_label:
                text_at(slide, x=526, y=360, w=360, h=18,
                        text=list_label, size=12,
                        color_hex=_EXEC["slate_light"])
            list_items = list(right.get("list") or [])
            if list_items:
                bullets_at(slide, x=526, y=384, w=360, h=100,
                           bullets=list_items, size=13,
                           color_hex=_EXEC["white"],
                           bullet_color_hex=_EXEC["cyan_bright"],
                           lh=1.7)

        elif layout == "exec_argument":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            bullets = list(raw.get("bullets") or [])
            takeaway = str(raw.get("takeaway") or "").strip()
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["blue"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=80, text=title_text,
                    size=36, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=152, w=880, h=36,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            add_rounded_rect(slide, x=48, y=210, w=864, h=270,
                             fill_hex=_EXEC["white"])
            bullets_at(slide, x=72, y=234, w=816, h=222,
                       bullets=bullets, size=15,
                       color_hex=_EXEC["ink"],
                       bullet_color_hex=_EXEC["blue"], lh=1.7)
            if takeaway:
                text_at(slide, x=48, y=496, w=864, h=28,
                        text=takeaway.upper(), size=13,
                        color_hex=_EXEC["blue"], bold=True,
                        align="center", anchor="middle")

        elif layout == "exec_stat_grid_4":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            stats = list(raw.get("stats") or [])
            detail_heading = str(raw.get("detail_heading") or "").strip()
            details = list(raw.get("details") or [])
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["blue"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=72, text=title_text,
                    size=36, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=142, w=880, h=30,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            xs = (48, 274, 500, 726)
            widths = (208, 208, 208, 186)
            accent_colors = (_EXEC["blue"], _EXEC["cyan"],
                             _EXEC["purple"], _EXEC["green"])
            for i in range(4):
                stat = stats[i] if i < len(stats) else {}
                if not isinstance(stat, Mapping):
                    stat = {}
                add_rect(slide, x=xs[i], y=200, w=widths[i], h=4,
                         fill_hex=accent_colors[i])
                add_rounded_rect(slide, x=xs[i], y=204, w=widths[i],
                                 h=144, fill_hex=_EXEC["white"])
                if stat.get("value") is not None:
                    text_at(slide, x=xs[i] + 20, y=222,
                            w=widths[i] - 40, h=70,
                            text=str(stat["value"]), size=48,
                            color_hex=_EXEC["ink"], bold=True,
                            font=t["font_head"])
                if stat.get("label"):
                    text_at(slide, x=xs[i] + 20, y=298,
                            w=widths[i] - 40, h=38,
                            text=str(stat["label"]), size=12,
                            color_hex=_EXEC["slate"])
            if detail_heading or details:
                if detail_heading:
                    text_at(slide, x=48, y=372, w=880, h=26,
                            text=detail_heading, size=15,
                            color_hex=_EXEC["ink"], bold=True)
                d_xs = (68, 354, 640)
                d_colors = (_EXEC["blue"], _EXEC["cyan"],
                            _EXEC["purple"])
                d_glyphs = ("◆", "◆", "◆")
                for i in range(3):
                    if i >= len(details):
                        break
                    d = details[i] if isinstance(details[i], Mapping) else {}
                    glyph = str(d.get("icon") or d_glyphs[i])
                    text_at(slide, x=d_xs[i], y=416, w=28, h=28,
                            text=glyph, size=20,
                            color_hex=d_colors[i], bold=True,
                            anchor="middle", align="center")
                    if d.get("title"):
                        text_at(slide, x=d_xs[i] + 42, y=414,
                                w=200, h=22,
                                text=str(d["title"]), size=13,
                                color_hex=_EXEC["ink"], bold=True)
                    if d.get("body"):
                        text_at(slide, x=d_xs[i] + 42, y=438,
                                w=200, h=60,
                                text=str(d["body"]), size=11,
                                color_hex=_EXEC["slate"])

        elif layout == "exec_features_2x2":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            features = list(raw.get("features") or [])
            banner_text = str(raw.get("banner_text") or "").strip()
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["purple"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=72, text=title_text,
                    size=36, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=142, w=880, h=30,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            f_xs = (48, 498, 48, 498)
            f_ys = (196, 196, 356, 356)
            f_ws = (430, 414, 430, 414)
            f_hs = (152, 152, 132, 132)
            ic_xs = (70, 520, 70, 520)
            ic_ys = (220, 220, 380, 380)
            t_xs = (150, 600, 150, 600)
            t_ws = (308, 292, 308, 292)
            colors = (_EXEC["blue"], _EXEC["cyan"],
                      _EXEC["purple"], _EXEC["green"])
            glyphs = ("✱", "◆", "▲", "●")
            for i in range(4):
                f = features[i] if i < len(features) else {}
                if not isinstance(f, Mapping):
                    f = {}
                add_rounded_rect(slide, x=f_xs[i], y=f_ys[i],
                                 w=f_ws[i], h=f_hs[i],
                                 fill_hex=_EXEC["white"])
                add_circle(slide, x=ic_xs[i], y=ic_ys[i], size=56,
                           fill_hex=colors[i])
                text_at(slide, x=ic_xs[i], y=ic_ys[i], w=56, h=56,
                        text=str(f.get("icon") or glyphs[i]),
                        size=24, color_hex=_EXEC["white"], bold=True,
                        align="center", anchor="middle")
                if f.get("title"):
                    text_at(slide, x=t_xs[i], y=ic_ys[i] + 6,
                            w=t_ws[i], h=28,
                            text=str(f["title"]), size=18,
                            color_hex=_EXEC["ink"], bold=True)
                if f.get("body"):
                    text_at(slide, x=t_xs[i], y=ic_ys[i] + 38,
                            w=t_ws[i], h=76,
                            text=str(f["body"]), size=12,
                            color_hex=_EXEC["slate"])
            if banner_text:
                add_rounded_rect(slide, x=48, y=502, w=864, h=26,
                                 fill_hex=_EXEC["navy_bg"])
                text_at(slide, x=48, y=502, w=864, h=26,
                        text=banner_text.upper(), size=12,
                        color_hex=_EXEC["cyan_bright"], bold=True,
                        align="center", anchor="middle")

        elif layout == "exec_industries_2x2":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            industries = list(raw.get("industries") or [])
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["blue"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=72, text=title_text,
                    size=36, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=142, w=880, h=30,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            rule_xs = (48, 488, 48, 488)
            bg_xs = (52, 492, 52, 492)
            bg_ws = (426, 420, 426, 420)
            ys = (196, 196, 360, 360)
            colors = (_EXEC["blue"], _EXEC["green"],
                      _EXEC["orange"], _EXEC["purple"])
            glyphs = ("◆", "▲", "●", "■")
            for i in range(4):
                ind = industries[i] if i < len(industries) else {}
                if not isinstance(ind, Mapping):
                    ind = {}
                add_rect(slide, x=rule_xs[i], y=ys[i], w=4, h=152,
                         fill_hex=colors[i])
                add_rounded_rect(slide, x=bg_xs[i], y=ys[i],
                                 w=bg_ws[i], h=152,
                                 fill_hex=_EXEC["white"])
                ic_x = bg_xs[i] + 24
                add_circle(slide, x=ic_x, y=ys[i] + 20, size=40,
                           fill_hex=colors[i])
                text_at(slide, x=ic_x, y=ys[i] + 20, w=40, h=40,
                        text=str(ind.get("icon") or glyphs[i]),
                        size=18, color_hex=_EXEC["white"], bold=True,
                        align="center", anchor="middle")
                name_x = ic_x + 54
                if ind.get("name"):
                    text_at(slide, x=name_x, y=ys[i] + 24, w=330, h=32,
                            text=str(ind["name"]), size=19,
                            color_hex=_EXEC["ink"], bold=True)
                # 4 check-marked uses, 2 columns × 2 rows.
                uses = list(ind.get("uses") or [])
                left_uses = uses[0::2][:2]
                right_uses = uses[1::2][:2]
                bullets_at(slide, x=bg_xs[i] + 24, y=ys[i] + 74,
                           w=(bg_ws[i] - 48) // 2,
                           h=72, bullets=left_uses, size=12,
                           color_hex=_EXEC["slate_dark"],
                           bullet_color_hex=_EXEC["green"],
                           style="check", lh=1.7)
                bullets_at(slide, x=bg_xs[i] + 24 + (bg_ws[i] - 48) // 2,
                           y=ys[i] + 74,
                           w=(bg_ws[i] - 48) // 2,
                           h=72, bullets=right_uses, size=12,
                           color_hex=_EXEC["slate_dark"],
                           bullet_color_hex=_EXEC["green"],
                           style="check", lh=1.7)

        elif layout == "exec_sovereignty_dark":
            paint_background(slide, _EXEC["navy_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            cards = list(raw.get("cards") or [])
            gov_heading = str(raw.get("gov_heading") or "").strip()
            governance = list(raw.get("governance") or [])
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["cyan_bright"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=72, text=title_text,
                    size=36, color_hex=_EXEC["white"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=142, w=880, h=30,
                        text=subhead, size=14,
                        color_hex=_EXEC["slate_light"])
            xs = (48, 274, 500, 726)
            widths = (208, 208, 208, 186)
            for i in range(4):
                c = cards[i] if i < len(cards) else {}
                if not isinstance(c, Mapping):
                    c = {}
                add_rounded_rect(slide, x=xs[i], y=196, w=widths[i],
                                 h=192, fill_hex=_EXEC["navy_card"],
                                 stroke_hex=_EXEC["navy_border"])
                if c.get("title"):
                    text_at(slide, x=xs[i] + 20, y=220,
                            w=widths[i] - 40, h=32,
                            text=str(c["title"]), size=22,
                            color_hex=_EXEC["cyan_bright"], bold=True,
                            font=t["font_head"])
                if c.get("body"):
                    text_at(slide, x=xs[i] + 20, y=280,
                            w=widths[i] - 40, h=96,
                            text=str(c["body"]), size=13,
                            color_hex=_EXEC["white"])
            if gov_heading or governance:
                add_rounded_rect(slide, x=48, y=408, w=864, h=96,
                                 fill_hex=_EXEC["light_bg"])
                if gov_heading:
                    text_at(slide, x=68, y=422, w=824, h=22,
                            text=gov_heading.upper(), size=12,
                            color_hex=_EXEC["blue"], bold=True)
                gov_xs = (68, 280, 492, 704)
                for i in range(4):
                    if i >= len(governance):
                        break
                    g = governance[i] if isinstance(
                        governance[i], Mapping) else {}
                    text_at(slide, x=gov_xs[i], y=450, w=16, h=20,
                            text="✓", size=12,
                            color_hex=_EXEC["green"], bold=True)
                    if g.get("title"):
                        text_at(slide, x=gov_xs[i] + 22, y=450,
                                w=180, h=20,
                                text=str(g["title"]), size=12,
                                color_hex=_EXEC["ink"], bold=True)
                    if g.get("body"):
                        text_at(slide, x=gov_xs[i] + 22, y=472,
                                w=180, h=24,
                                text=str(g["body"]), size=10,
                                color_hex=_EXEC["slate"])

        elif layout == "exec_chat_example":
            paint_background(slide, _EXEC["light_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            subhead = str(raw.get("subhead") or "").strip()
            chat = raw.get("chat") or {}
            if not isinstance(chat, Mapping):
                chat = {}
            stats = list(raw.get("stats") or [])
            if kicker:
                text_at(slide, x=48, y=36, w=700, h=20,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["cyan"], bold=True)
            text_at(slide, x=48, y=62, w=880, h=72, text=title_text,
                    size=36, color_hex=_EXEC["ink"], bold=True,
                    font=t["font_head"])
            if subhead:
                text_at(slide, x=48, y=142, w=880, h=30,
                        text=subhead, size=14, color_hex=_EXEC["slate"])
            # Left chat panel.
            add_rounded_rect(slide, x=48, y=200, w=510, h=308,
                             fill_hex=_EXEC["chat_panel"])
            if chat.get("kicker"):
                text_at(slide, x=72, y=220, w=460, h=18,
                        text=str(chat["kicker"]).upper(), size=11,
                        color_hex=_EXEC["slate"], bold=True)
            # User bubble (blue).
            add_rounded_rect(slide, x=92, y=252, w=446, h=50,
                             fill_hex=_EXEC["blue"])
            if chat.get("question"):
                text_at(slide, x=108, y=252, w=414, h=50,
                        text=str(chat["question"]), size=12,
                        color_hex=_EXEC["white"], anchor="middle")
            # Answer card (white).
            add_rounded_rect(slide, x=72, y=314, w=462, h=144,
                             fill_hex=_EXEC["white"])
            if chat.get("answer_headline"):
                text_at(slide, x=92, y=330, w=422, h=22,
                        text=str(chat["answer_headline"]), size=13,
                        color_hex=_EXEC["ink"], bold=True)
            if chat.get("answer_subline"):
                text_at(slide, x=92, y=354, w=422, h=20,
                        text=str(chat["answer_subline"]), size=11,
                        color_hex=_EXEC["slate"])
            answer_bullets = list(chat.get("answer_bullets") or [])
            if answer_bullets:
                bullets_at(slide, x=92, y=382, w=422, h=68,
                           bullets=answer_bullets, size=11,
                           color_hex=_EXEC["slate_dark"],
                           bullet_color_hex=_EXEC["slate"], lh=1.6)
            if chat.get("source"):
                add_rounded_rect(slide, x=72, y=472, w=220, h=22,
                                 fill_hex=_EXEC["green_soft"])
                text_at(slide, x=72, y=472, w=220, h=22,
                        text=str(chat["source"]), size=10,
                        color_hex=_EXEC["green_ink"], bold=True,
                        align="center", anchor="middle")
            # Right: 3 stat blocks.
            stat_ys = (200, 308, 416)
            for i in range(3):
                if i >= len(stats):
                    break
                s = stats[i] if isinstance(stats[i], Mapping) else {}
                add_rect(slide, x=580, y=stat_ys[i], w=4, h=90,
                         fill_hex=_EXEC["cyan"])
                if s.get("value") is not None:
                    text_at(slide, x=596, y=stat_ys[i],
                            w=200, h=50,
                            text=str(s["value"]), size=38,
                            color_hex=_EXEC["ink"], bold=True,
                            font=t["font_head"])
                if s.get("label"):
                    text_at(slide, x=596, y=stat_ys[i] + 56,
                            w=316, h=20,
                            text=str(s["label"]), size=13,
                            color_hex=_EXEC["ink"], bold=True)
                if s.get("sub"):
                    text_at(slide, x=596, y=stat_ys[i] + 78,
                            w=316, h=18,
                            text=str(s["sub"]), size=11,
                            color_hex=_EXEC["slate"])

        elif layout == "exec_closing_dark":
            paint_background(slide, _EXEC["navy_bg"])
            kicker = str(raw.get("kicker") or "").strip()
            title_text = str(raw.get("title") or "")
            reasons = list(raw.get("reasons") or [])
            cta = str(raw.get("cta") or "").strip()
            if kicker:
                text_at(slide, x=48, y=56, w=600, h=22,
                        text=kicker.upper(), size=12,
                        color_hex=_EXEC["cyan_bright"], bold=True)
            text_at(slide, x=48, y=92, w=880, h=68, text=title_text,
                    size=44, color_hex=_EXEC["white"], bold=True,
                    font=t["font_head"])
            xs = (48, 502, 48, 502)
            ys = (196, 196, 356, 356)
            widths = (430, 410, 430, 410)
            for i in range(4):
                r = reasons[i] if i < len(reasons) else {}
                if not isinstance(r, Mapping):
                    r = {}
                add_rounded_rect(slide, x=xs[i], y=ys[i],
                                 w=widths[i], h=144,
                                 fill_hex=_EXEC["navy_card"],
                                 stroke_hex=_EXEC["navy_border"])
                num = r.get("number") or f"{i + 1:02d}"
                text_at(slide, x=xs[i] + 22, y=ys[i] + 16,
                        w=100, h=66,
                        text=str(num), size=44,
                        color_hex=_EXEC["cyan_bright"], bold=True,
                        font=t["font_head"])
                if r.get("title"):
                    text_at(slide, x=xs[i] + 130, y=ys[i] + 24,
                            w=widths[i] - 150, h=28,
                            text=str(r["title"]), size=18,
                            color_hex=_EXEC["white"], bold=True)
                if r.get("body"):
                    text_at(slide, x=xs[i] + 130, y=ys[i] + 56,
                            w=widths[i] - 150, h=80,
                            text=str(r["body"]), size=12,
                            color_hex=_EXEC["slate_light"])
            if cta:
                add_rounded_rect(slide, x=48, y=504, w=864, h=24,
                                 fill_hex=_EXEC["cyan_bright"])
                text_at(slide, x=48, y=504, w=864, h=24,
                        text=cta.upper(), size=13,
                        color_hex=_EXEC["navy_bg"], bold=True,
                        align="center", anchor="middle")

        # Footer page number, except on title / section_divider / exec_*.
        if (layout not in ("title", "section_divider", "full_bleed_image")
                and layout not in _EXEC_LAYOUTS):
            pn = slide.shapes.add_textbox(
                SW - Inches(1.0), SH - Inches(0.45),
                Inches(0.6), Inches(0.3),
            )
            set_text(pn.text_frame, str(idx),
                     size=t["size_kicker"], color_hex=t["text_muted"],
                     align=PP_ALIGN.RIGHT)

    out = _stamped_path(title, "pptx")
    prs.save(str(out))
    _maybe_publish(out, kind="pptx", title=title, summary=summary,
                   auto_publish=auto_publish)
    return out


# ------------------------------- Chart -------------------------------------
def make_chart_png(
    spec: Mapping[str, Any],
    *,
    title: str = "chart",
    summary: str = "",
    auto_publish: bool = True,
) -> Path:
    """Render a simple chart from a tiny spec dict.

    Spec shape::

        {
          "kind": "bar" | "line" | "scatter",
          "x": [...], "y": [...],
          "x_label": "...", "y_label": "...", "title": "..."
        }

    Anything more complex — write your own matplotlib code and save into
    ``citra_toolkit.report._reports_dir()`` directly, then ``publish()``.
    """
    import matplotlib  # type: ignore[import-not-found]

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # type: ignore[import-not-found]

    kind = (spec.get("kind") or "bar").lower()
    x = list(spec.get("x") or [])
    y = list(spec.get("y") or [])
    fig, ax = plt.subplots(figsize=(8, 4.5), dpi=120)
    if kind == "line":
        ax.plot(x, y)
    elif kind == "scatter":
        ax.scatter(x, y)
    else:
        ax.bar([str(v) for v in x], y)
    if spec.get("title"):
        ax.set_title(str(spec["title"]))
    if spec.get("x_label"):
        ax.set_xlabel(str(spec["x_label"]))
    if spec.get("y_label"):
        ax.set_ylabel(str(spec["y_label"]))
    fig.autofmt_xdate(rotation=30)
    fig.tight_layout()
    out = _stamped_path(title, "png")
    fig.savefig(str(out))
    plt.close(fig)
    _maybe_publish(out, kind="png", title=title, summary=summary, auto_publish=auto_publish)
    return out


# Path to the baked-in Chart.js renderer (see infrastructure/action-sandbox/
# chartjs/ and the Dockerfile). Overridable for tests / non-standard layouts.
_CHARTJS_RENDERER = os.getenv("CITRA_CHARTJS_RENDERER", "/srv/citra/chartjs/render.js")


def _chartjs_render(
    config: Mapping[str, Any],
    *,
    width: int = 900,
    height: int = 500,
    background: str = "white",
) -> bytes:
    """Render a Chart.js config to raw PNG bytes via the baked-in renderer.

    Shells out to ``node /srv/citra/chartjs/render.js`` (``chart.js`` +
    ``chartjs-node-canvas``, baked into the sandbox image). Used by
    :func:`make_chartjs_png` and by the ``chart`` element of a freeform
    ``make_pptx`` slide.
    """
    import json
    import subprocess

    payload = json.dumps(
        {
            "config": dict(config),
            "width": int(width),
            "height": int(height),
            "background": background,
        }
    )

    try:
        proc = subprocess.run(
            ["node", _CHARTJS_RENDERER],
            input=payload.encode("utf-8"),
            capture_output=True,
            timeout=60,
        )
    except FileNotFoundError as exc:  # `node` missing — not a sandbox image
        raise RuntimeError(
            "Chart.js renderer unavailable: `node` not found. "
            "Chart.js rendering only works inside the sandbox image."
        ) from exc

    if proc.returncode != 0 or not proc.stdout:
        detail = proc.stderr.decode("utf-8", "replace").strip()[:500]
        raise RuntimeError(f"Chart.js render failed: {detail or 'no output'}")

    return proc.stdout


def make_chartjs_png(
    config: Mapping[str, Any],
    *,
    title: str = "chart",
    summary: str = "",
    width: int = 900,
    height: int = 500,
    background: str = "white",
    auto_publish: bool = True,
) -> Path:
    """Render a Chart.js config to a PNG via the baked-in Node renderer.

    ``config`` is a standard Chart.js configuration object — exactly what
    you would pass to ``new Chart(ctx, config)`` in a browser::

        make_chartjs_png({
            "type": "line",
            "data": {
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "datasets": [{"label": "FY26", "data": [2.8, 3.0, 3.2, 3.4]}],
            },
            "options": {"plugins": {"title": {"display": True, "text": "Coverage"}}},
        })

    Chart.js is JavaScript and cannot run inside WeasyPrint (PDF) or
    python-pptx (PPTX), so the chart is pre-rendered to a PNG here:

    - **PPTX** — pass the returned ``Path`` in a slide image element, or
      skip this helper entirely and use a ``chart`` element on a freeform
      slide (``make_pptx`` renders the Chart.js config inline).
    - **PDF** — base64-encode the bytes into an ``<img>`` tag in the HTML
      you hand to :func:`make_pdf`.

    For in-process analytical plots from a DataFrame, prefer matplotlib
    via :func:`make_chart_png`.
    """
    png = _chartjs_render(config, width=width, height=height, background=background)
    out = _stamped_path(title, "png")
    out.write_bytes(png)
    _maybe_publish(out, kind="png", title=title, summary=summary, auto_publish=auto_publish)
    return out
